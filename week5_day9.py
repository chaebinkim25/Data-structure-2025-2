from pptx import __version__ as pptx_version
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

print(pptx_version)

# 새 프레젠테이션 생성
prs = Presentation("template.pptx")
prs.slides.add_slide(prs.slide_masters[0].slide_layouts[0])  # 제목 슬라이드 추가
# 제목 슬라이드 설정
title_slide = prs.slides[0]
title_slide.shapes.title.text = "자료구조 (Data Structure)\n5주차: 트리"


layout = prs.slide_masters[0].slide_layouts[1]  # 제목과 내용 레이아웃 선택

# 슬라이드 추가 함수
def add_slide(title, bullets):
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    content = slide.placeholders[1]
    content.text = bullets[0] if bullets else ""
    for bullet in bullets[1:]:
        p = content.text_frame.add_paragraph()
        p.text = bullet

slides = [
    ("선형 (일차원) 자료구조 정리", [
        "리스트: 데이터를 추가하면 순서가 유지되는 구조","",
        "스택: 가장 최신 데이터만 꺼낼 수 있는 구조","",
        "큐: 가장 오래된 데이터만 꺼낼 수 있는 구조","",
    ]),
    ("이번 시간 목차", [
        "새로운 자료구조 소개","",
        "새로운 자료구조 예제","",
        "새로운 자료구조의 데이터들에서 반복하는 방법","",
    ]),
    ("새로운 자료 구조", [
        "리스트의 노드마다 리스트를 가지고 있다면?","",
        "예시) 파일 시스템, 소스 코드","",
    ]),
    ("새로운 자료 구조", [
        "세부 폴더들은 폴더에 연결. 일대다 (부모 자식) 관계","",
        "모든 폴더의 조상 (루트) 존재","",
    ]),
    ("트리란?", [
        "트리 정의: 루트로부터 부모 자식 관계로 연결된 구조","",
    ]),
    ("이진 트리", [
        "자식 노드가 최대 2개인 트리","",
        "\t - root: 트리의 시작점 (타입: 이진 트리 노드)","",
    ]),
    ("이진 트리 노드 타입", [
        "데이터와 자식 노드 2개 정보의 묶음","",
        "\t - left: 왼쪽 자식 노드 주소","",
        "\t - right: 오른쪽 자식 노드 주소","",
        "\t - value: 저장된 데이터","",
    ]),
    ("이진 트리 예시", [
        "다음과 같은 이진 트리를 만들어보자","",
        "트리 T를 만든다",
        "노드 A, B, C, D, E를 만든다",
        "T의 root = A",
        "A의 left = B",
        "A의 right = C",
        "B의 left = D",
        "B의 right = E",
    ]),
    ("예제: 더하기/곱하기 수식 트리", [
        "1 + 2 + 3 수식을 이진 트리로 만들면 계산 순서가 명확히 드러난다","",
        "트리 T를 만든다",
        "노드 1을 만든다",
        "T의 root = 노드 1",
        "노드 +를 만든다",
        "노드 +의 left = 노드 1",
        "T의 root = 노드 +",
        "노드 2를 만든다",
        "노드 +의 right = 노드 2",
        "두번째 + 노드를 만든다",
        "두번째 + 노드의 left = 첫번째 + 노드",
        "T의 root = 두번째 + 노드",
        "노드 3을 만든다",
        "두번째 + 노드의 right = 노드 3",
    ]),
    ("예제: 더하기/곱하기 수식 트리", [
        "1 + 2 * 3 수식을 이진 트리로 만들면 계산 순서가 명확히 드러난다","",
        "트리 T를 만든다",
        "노드 1을 만든다",
        "T의 root = 노드 1",
        "노드 +를 만든다",
        "노드 +의 left = 노드 1",
        "T의 root = 노드 +",
        "노드 2를 만든다",
        "노드 +의 right = 노드 2",
        "노드 *을 만든다",
        "노드 *의 left = 노드 2",
        "노드 +의 right = 노드 *",
        "노드 3을 만든다",
        "노드 *의 right = 노드 3",
    ]),
    ("예제: 더하기/곱하기 수식 토큰", [
        "숫자 또는 연산자가 저장된다","",
        "\t - type: 숫자, 더하기, 곱하기 구분",
        "\t - value: 숫자일 경우 값을 저장",
        "\t - precedence: 연산자 우선순위",
        "\t               더하기는 1",
        "\t               곱하기는 2",
        "\t               숫자는 -1 (연산자 아님)",
    ]),
    ("예제: 더하기/곱하기 수식 트리 생성", [
        "함수 expr_tree(수식 토큰 큐 Q, 최소 우선순위 min_prec):",
        "node = 새 노드 (dequeue(Q), NULL, NULL)",
        "반복: size(Q) > 0",
        "\top = front(Q)",
        "\t만약: 우선순위(op) < min_prec이면 반복 종료",
        "\tdequeue(Q)",
        "\tnode = 새 노드 (op, node, expr_tree(Q, 우선순위(op) + 1))",
        "반환: node",
    ]),
    ("예제: 더하기/곱하기 수식 트리 생성", [
        "expr_tree([1, '+', 2, '+', 3], 0)의 실행 과정은?","",
        "1. node = 새 노드 (1, NULL, NULL)\n",
        "2. op = '+'\n",
        "3. node = 새 노드 ('+', node, expr_tree([2, '+', 3]))\n",
        "\t  3-1. node = 새 노드 (2, NULL, NULL)\n",
        "\t  3-2. op = '+'\n",
        "\t  3-3. 반환: node",
        "4. 다음 반복",
        "5. op = '+'\n",
        "6. node = 새 노드 ('+', node, expr_tree([3]))\n",
        "\t  6-1. node = 새 노드 (3, NULL, NULL)\n",
        "\t  6-2. 반환: node",
        "7. 다음 반복",
        "8. 반복 종료",
        "9. 반환: node",
    ]),
    ("예제: 더하기/곱하기 수식 트리 생성", [
        "expr_tree([1, '+', 2, '*', 3], 0)의 실행 과정은?","",
        "1. node = 새 노드 (1, NULL, NULL)\n",
        "2. op = '+'\n",
        "3. node = 새 노드 ('+', node, expr_tree([2, '*', 3], 2))\n",
        "\t  3-1. node = 새 노드 (2, NULL, NULL)\n",
        "\t  3-2. op = '*'\n",
        "\t  3-3. node = 새 노드 ('*', node, expr_tree([3], 3))\n",
        "\t      3-3-1. node = 새 노드 (3, NULL, NULL)\n",
        "\t      3-3-2. 반환: node\n",
        "\t  3-4. 반환: node",
        "4. 다음 반복",
        "5. 반복 종료",
        "6. 반환: node",
    ]),

    ("예제: 더하기/곱하기 수식 트리 계산", [
        "함수 EVAL(node):",
        "\t만약: node가 숫자 노드면",
        "\t\t반환: node의 숫자 값",
        "\t만약: node가 연산자 노드면",
        "\t\tleft_val = EVAL(node->left)",
        "\t\tright_val = EVAL(node->right)",
        "\t\t만약: node가 '+'면 ",
        "\t\t\t반환: left_val + right_val",
        "\t\t만약: node가 '*'이면 ",
        "\t\t\t반환: left_val * right_val",
    ]),
    ("예제: 더하기/곱하기 수식 트리 계산", [
        "'1 + 2 + 3' 수식에서 EVAL(root)의 실행 과정은?","",
        "1. node '+'는 연산자 노드",
        "2. left_val = EVAL(node->left)",
        "\t  2-1. node '+'는 연산자 노드",
        "\t  2-2. left_val = EVAL(node->left)",
        "\t      2-2-1. node '1'은 숫자 노드",
        "\t      2-2-2. 반환: 1",
        "\t  2-3. right_val = EVAL(node->right)",
        "\t      2-3-1. node '2'는 숫자 노드",
        "\t      2-3-2. 반환: 2",
        "\t  2-4. 반환: 1 + 2",
        "3. right_val = EVAL(node->right)",
        "\t  3-1. node '3'은 숫자 노드",
        "\t  3-2. 반환: 3",
        "4. 반환: 3 + 3",
    ]),
    ("예제: 더하기/곱하기 수식 트리 계산", [
        "'1 + 2 * 3' 수식에서 EVAL(root)의 실행 과정은?","",
        "1. node '+'는 연산자 노드",
        "2. left_val = EVAL(node->left)",
        "\t  2-1. node '1'은 숫자 노드",
        "\t  2-2. 반환: 1",
        "3. right_val = EVAL(node->right)",
        "\t  3-1. node '*'는 연산자 노드",
        "\t  3-2. left_val = EVAL(node->left)",
        "\t      3-2-1. node '2'는 숫자 노드",
        "\t      3-2-2. 반환: 2",
        "\t  3-3. right_val = EVAL(node->right)",
        "\t      3-3-1. node '3'는 숫자 노드",
        "\t      3-3-2. 반환: 3",
        "\t  3-4. 반환: 2 * 3",
        "4. 반환: 1 + 6",
    ]),
    ("이진 트리 방문", [
        "목적: 이진 트리의 모든 노드에 대해 반복하기","",
        "깊이 우선: 자손들 방문이 끝나고 다음 반복 진행",
        "\t 데이터 처리 타이밍:",
        "\t - Preorder: 방문전 처리할 작업",
        "\t - Inorder: 방문 도중 처리할 작업",
        "\t - Postorder: 방문 후 처리할 작업",
        "너비 우선: 한 세대의 방문이 끝나고 다음 세대 방문",
        "\t - Level-order: 방문시 처리할 작업",
    ]),
    ("깊이 우선 이진 트리 방문 구현", [
        "재귀 함수 호출:",
        "\t 한 노드에 대한 동작만 정하면,",
        "\t 노드들 전체에 적용된다",
        "",
        "스택:\t 최적화가 가능하다",
    ]),
    ("이진 트리 방문시 작업 타이밍 - Preorder", [
        "데이터 전처리를 먼저 하고, 자식들에서 반복한다","",
        "함수 TRAVERSE_PRE(node):",
        "\t만약: node == NULL이면 반환",
        "\tpre_task(node->value)",
        "\tTRAVERSE_PRE(node->left)",
        "\tTRAVERSE_PRE(node->right)",
    ]),
    ("이진 트리 방문시 작업 타이밍 - Preorder", [
        "다음 트리에 대한 TRAVERSE_PRE(root)에서",
        "전처리 작업이 실행되는 순서는?",
        "",
        "       A",
        "      / \\",
        "     B   C",
        "    / \\",
        "   D   E",
        "",
        "A → B → D → E → C",
    ]),
    ("이진 트리 방문시 작업 타이밍 - Inorder", [
        "왼쪽 자손들 방문이 끝나면 데이터 중간 처리를 한다.","",
        "함수 TRAVERSE_IN(node):",
        "\t만약: node == NULL이면 반환",
        "\tTRAVERSE_IN(node->left)",
        "\tmiddle_task(node->value)",
        "\tTRAVERSE_IN(node->right)",
    ]),
    ("이진 트리 방문시 작업 타이밍 - Inorder", [
        "다음 트리에 대한 TRAVERSE_IN(root)에서",
        "중간 처리 작업이 실행되는 순서는?",
        "",
        "       A",
        "      / \\",
        "     B   C",
        "    / \\",
        "   D   E",
        "",
        "D → B → E → A → C",
    ]),
    ("이진 트리 방문시 작업 타이밍 - Postorder", [
        "모든 자손들 방문이 끝나면 데이터 후처리를 한다.","",
        "함수 TRAVERSE_POST(node):",
        "\t만약: node == NULL이면 반환",
        "\tTRAVERSE_POST(node->left)",
        "\tTRAVERSE_POST(node->right)",
        "\tpost_task(node->value)",
    ]),
    ("이진 트리 방문시 작업 타이밍 - Postorder", [
        "다음 트리에 대한 TRAVERSE_POST(root)에서",
        "후처리 작업이 실행되는 순서는?",
        "",
        "       A",
        "      / \\",
        "     B   C",
        "    / \\",
        "   D   E",
        "",
        "D → E → B → C → A",
    ]),

    ("이진 트리 방문 퀴즈", [
        "Preorder traversal\tA B C D E F G",
        "Inorder traversal\tC D B A E G F",
        "Original tree???",
        "",
        "       A\n",
        "      / \\\n",
        "     B   E\n",
        "    /     \\\n",
        "   C      F\n",
        "    \\   /\n",
        "     D  G\n",
        "",
    ]),

    ("요약", [
        "트리: 루트로부터 부모 자식 관계로 연결된 구조","",
        "이진 트리: 자식 노드가 최대 2개인 트리","",
        "이진 트리 노드: left, right, value","",
        "트리 방문: Preorder, Inorder, Postorder","",
        "다음 시간: C로 구현하기","",
    ]),

    ("이진 트리 방문 - 재귀함수로 구현", [
        "자손들 방문은 왼쪽, 오른쪽 순으로 한다. ",
        "데이터는 전처리, 중간처리, 후처리를 한다.",
        "",
        "함수 TRAVERSE(node):",
        "\t만약: node == NULL이면 반환",
        "\tpre_task(node->value)",
        "\tTRAVERSE(node->left)",
        "\tmiddle_task(node->value)",
        "\tTRAVERSE(node->right)",
        "\tpost_task(node->value)",
    ]),

    ("이진 트리 방문 - 노드 상태 구조체", [
        "최적화를 하기 위해 스택으로 구현한다","",
        "구조체 노드상태"
        "\t - node: 이진트리노드",
        "\t - state: 방문 상태",
        "",
        "방문 상태: ",
        "\t - 0(pre_task 전)",
        "\t - 1(middle_task 전)",
        "\t - 2(post_task 전)",
        "",
    ]),

    ("이진 트리 방문 - 스택으로 구현", [
        "함수 traverse(root):",
        "to_visit = 새로운 노드상태 스택",
        "PUSH(to_visit, 새 노드상태(root, 0))",
        "반복: to_visit이 비어있지 않으면",
        "\tns = PEEK(to_visit)",
        "\t만약: ns->node == NULL이면 다음 반복으로",
        "\t만약: ns->state == 0이면",
        "\t\tpre_task(ns->node->value)",
        "\t\tns->state = 1",
        "\t\tPUSH(to_visit, 새 노드상태(ns->node->left, 0))",
        "\t그 외에 만약: ns->state == 1이면",
        "\t\tmiddle_task(ns->node->value)",
        "\t\tns->state = 2",
        "\t\tPUSH(to_visit, 새 노드상태(ns->node->right, 0))",
        "\t그 외에 만약: ns->state == 2이면",
        "\t\tpost_task(ns->node->value)",
        "\t\tPOP(to_visit)",
        "",
    ]),

    ("이진 트리 방문 - 간소화", [
        "함수 traverse(root):",
        "parents = 새로운 노드 스택",
        "next_child = root",
        "last_child_done = NULL",
        "반복: next_child가 있거나 parents가 비어있지 않으면",
        "\t만약: next_child가 있으면",
        "\t\tpre_task(next_child->value)",
        "\t\tPUSH(parents, next_child)",
        "\t\tnext_child = next_child->left",
        "\t\t다음 반복으로",
        "\tparent = PEEK(parents)",
        "\t만약: parent->right == NULL",
        "\t\tPOP(parents)",
        "\t\tmiddle_task(parent->value)",
        "\t\tpost_task(parent->value)",
        "\t\tlast_child_done = parent",
        "\t\t다음 반복으로",
        "\t만약: last_child_done == parent->right",
        "\t\tPOP(parents)",
        "\t\tpost_task(parent->value)",
        "\t\tlast_child_done = parent",
        "\t\t다음 반복으로",
        "\tmiddle_task(parent->value)",
        "\tnext_child = parent->right",
    ]),


    ("이진 트리 방문 - Levelorder", [
        "함수 level_order(root)는 루트에서 가까운 순으로 방문한다","",
        "level_order(root):",
        "    q = 새로운 노드 큐",
        "    ENQUEUE(q, root)",
        "    반복: q가 비어있지 않으면",
        "        n = DEQUEUE(q)",
        "        만약: n이 있으면",
        "            level_task(n->value)",
        "            ENQUEUE(q, n->left)",
        "            ENQUEUE(q, n->right)",
    ]),

]

# 슬라이드 생성
for title, bullets in slides:
    add_slide(title, bullets)


def add_table_slides(title, bullets, table_data):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title

        content = slide.placeholders[1]
        content.text = bullets[0] if bullets else ""
        for bullet in bullets[1:]:
                p = content.text_frame.add_paragraph()
                p.text = bullet

        cols, rows = len(table_data), len(table_data[0]) if table_data else 0
        left = Inches(0.5)
        top = Inches(5)
        width = Inches(9)
        height = Inches(0.8 + 0.3 * rows)
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        for c in range(cols):
                for r in range(rows):
                        table.cell(r, c).text = table_data[c][r]


table_slides = [
        ("트리 방문 - Preorder", [
                "1    만약: N == NULL이면 반환",
                "2    pre_task(N)",
                "3    PREORDER(LEFT(N))",
                "4    PREORDER(RIGHT(N))",
        ], [
                ["호출스택", "pre_task", "줄번호"],
                ["A", "A", "3"],
                ["AB", "B", "3"],
                ["ABn", "", "1"],
                ["AB", "", "4"],
                ["AB", "", "-"],
                ["A", "", "4"],
                ["AC", "C", "3"],
                ["ACn", "", "1"],
                ["AC", "", "4"],
                ["ACn", "", "1"],
                ["AC", "", "4"],
                ["A", "", "4"],
                ["An", "", "1"],
        ]),
]

for title, bullets, table_data in table_slides:
        add_table_slides(title, bullets, table_data)


# 저장
output_path = "ppt/자료구조_5주차_9회차_PPT.pptx"
prs.save(output_path)
output_path
