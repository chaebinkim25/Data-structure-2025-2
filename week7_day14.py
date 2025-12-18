from pptx import __version__ as pptx_version
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

print(pptx_version)

# 새 프레젠테이션 생성
prs = Presentation("template.pptx")
prs.slides.add_slide(prs.slide_masters[0].slide_layouts[0])  # 제목 슬라이드 추가
# 제목 슬라이드 설정
title_slide = prs.slides[0]
title_slide.shapes.title.text = "자료구조 (Data Structure)\n7주차: 정렬"


layout = prs.slide_masters[0].slide_layouts[1]  # 제목과 내용 레이아웃 선택

# 슬라이드 추가 함수
def add_slide(title, bullets):
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    content = slide.placeholders[1]
    content.text = bullets[0] if bullets else ""
    for bullet in bullets[1:]:
        p = content.text_frame.add_paragraph()
        p.text = bullet

slides = [
    ("지난 시간 요약", [
        "heap sort: 시간 O(n log n)",
        "\t\t공간 O(1)","",
        "quick sort: 시간 평균 O(n log n) 최악 O(n²)",
        "\t\t공간 O(log n)","",
        "insertion sort: 시간 최선 O(n) 평균 O(n²)",
        "\t\t공간 O(1)",
    ]),
    ("오늘의 목차", [
        "배열에 저장된 트리 출력하기 (C 코드)","",
        "heap sort (리눅스 커널의 최초 버전)","",
    ]),
    ("트리 출력 - printf", [
        "#include <stdio.h>",
        "int main()",
        "{",
        '\tprintf("Hello");',
        "\treturn 0;",
        "}"
    ]),
    ("트리 출력 - printf", [
        "int printf(const char *format, ...);",
        "char *는 글자 값이 저장된 주소",
        "const char *는 해당 변수를 통해서 글자값을 바꾸지 않겠다는 약속",
    ]),
    ("트리 출력 - printf - 글자 배열", [
        "#include <stdio.h>",
        "int main()",
        "{",
        "\tchar arr[5];",
        "\tarr[0] = 'H';",
        "\tarr[1] = 'i';"
        "\tarr[2] = '\\0';",
        "\tprintf(arr);",
        "\treturn 0;",
        "}"
    ]),
    ("트리 출력 - printf - 글자 배열", [
        "글자에 해당하는 값은 작은 따옴표로 표시한다.",
        "글자 배열을 문장으로 다룰 때, 마지막 글자는 '\\0'라는 특수문자로 쓴다.",
    ]),
    ("트리 출력 - printf - 글자 배열에 초기값 저장", [
        "#include <stdio.h>",
        "int main()",
        "{",
        "\tchar arr[5] = {'H', 'i', '\\0'};",
        "\tprintf(arr);",
        "\treturn 0;",
        "}"
    ]),
    ("트리 출력 - printf - 글자 배열에 초기값 저장", [
        "#include <stdio.h>",
        "int main()",
        "{",
        '\tchar arr[5] = "Hi";',
        "\tprintf(arr);",
        "\treturn 0;",
        "}",
    ]),
    ("트리 출력 - printf - 글자 배열에 초기값 저장", [
        "배열을 처음 만들 때, 초기값들을 중괄호 { } 안에 적어줄 수 있다.",
        '문자 배열의 경우, 초기값들을 큰 따옴표"" 안에 적어줄 수 있다. 공간이 충분하면 마지막에 \'\\0\'도 붙는다.'
    ]),
    ("트리 출력 - printf - 특수문자 '\\0'", [
        "#include <stdio.h>",
        "int main()",
        "{",
        '\tchar arr[5] = {"H\\0i"};',
        "\tprintf(arr);",
        "\treturn 0;",
        "}"
    ]),
    ("트리 출력 - printf - 특수문자 '\\0'", [
        "글자 배열에서 글자를 하나씩 처리하다가 '\\0'을 발견하면 문장이 끝났다고 보고 처리를 종료한다.",
    ]),
    ("트리 출력 - printf - 포맷문자열 %c, %d", [
        "#include <stdio.h>",
        "int main()",
        "{",
        '\tchar arr[5] = {"H\\0i"};',
        "\tfor (int i = 0; i < 4; i++)",
        '\t\tprintf("%c[%d]\n", arr[i], arr[i]);',
        "\treturn 0;",
        "}"
    ]),
    ("트리 출력 - printf - 포맷문자열 %c, %d", [
        "%로 표시된 자리에는 다음 매개변수의 값을 적절한 서식에 맞춰 대입한다.",
        "%c는 문자 값을 문자로 나타낸다.",
        "%d는 숫자 값을 숫자로 나타낸다.",
        "\\n은 줄바꿈 특수문자를 의미한다.",
    ]),
    ("트리 출력 - printf - 포맷문자열의 서식 너비", [
        "#include <stdio.h>",
        "int main()",
        "{",
        '\tprintf("%5c\n", \'A\');',    
        '\tprintf("%4c\n", \'B\');',   
        '\tprintf("%3c\n", \'C\');',   
        '\tprintf("%2c\n", \'D\');',   
        '\tprintf("%1c\n", \'E\');',   
        "\treturn 0;",
        "}"
    ]),
    ("트리 출력 - printf - 포맷문자열의 서식 너비", [
        "%d와 서식 문자 사이에 숫자를 적으면, 해당 항목의 최소 너비를 정할 수 있다.",
    ]),
    ("트리 출력 - printf - 포맷문자열의 가변 서식 너비", [
        "#include <stdio.h>",
        "int main()",
        "{",
        '\tfor (int i = 1; i <= 5; i++)',    
        '\t\tprintf("%*c\n", \'z\');',   
        "\treturn 0;",
        "}",
    ]),
    ("트리 출력 - printf - 포맷문자열의 가변 서식 너비", [
        "%d와 서식 문자 사이에 *을 적으면, 해당 항목의 최소 너비를 매개변수로 전달할 수 있다.",
    ]),
    ("트리 출력 - 계획", [
        "---------------32---------------",
        "-------16-------A------16-------",
        "---8----B------16-------C---8---",
        "-4-D----8---E---8---F---8---G-4-",
    ]),
    ("트리 출력 - 계획", [
        "#include <stdio.h>",
        "int main()",
        "{",
        '\tprintf("%16c\n", \'A\');',"",
        '\tprintf("%8c", \'B\');',
        '\tprintf("%16c\n", \'C\');',"",
        '\tprintf("%4c", \'D\');',
        '\tprintf("%8c", \'E\');',
        '\tprintf("%8c", \'F\');',
        '\tprintf("%8c\n", \'G\');',
        "\treturn 0;",
        "}",
    ]),
    ("트리 출력 - 계획", [
        "#include <stdio.h>",
        "int main()",
        "{",
        '\tint total_width = 32;',
        '\tchar data[10] = "ABCDEFG";',
        '\tint level = 1;',
        '\tint level_size = 1;',
        '\tint level_width = total_width;',
        '\tint level_start = 0;',
        '\tint level_end = 0;',
        '\tfor (int i = 0; i < 7; i++) {',
        '\t\tint spaces = level_width;',
        '\t\tif (i == level_start)',
        '\t\t\tspaces /= 2;',
        '\t\tprintf("%*c", spaces, data[i]);',
        '\t\tif (i == level_end) {',
        '\t\t\tprintf("\n");',
        '\t\t\tlevel += 1;',
        '\t\t\tlevel_size *= 2;',
        '\t\t\tlevel_width = total_width / level_size;',
        '\t\t\tlevel_start = level_end + 1;',
        '\t\t\tlevel_end += level_size;',
        '\t\t}',
        '\t}',
        '}',
    ]),
    ("트리 출력 - print_tree()", [
        "void print_tree(char *data, int tree_size, int tree_width)",
        "{",
        "\tint level_size = 1;",
        "\tint node_width = tree_width;",
        "\tint level_start = 0;", 
        "\tint level_end = 0;",
        "\tfor (int i = 0; i < tree_size; i++) {",
        "\t\tint spaces = node_width;",
        "\t\tif (i == level_start)",
        "\t\t\tspaces /= 2;",
        '\t\tprintf("%*c", spaces, data[i]);',
        "\t\tif (i == level_end) {",
        '\t\t\tprintf("\\n");',
        "\t\t\tlevel_size *= 2;",
        "\t\t\tnode_width = tree_width / level_size;",
        '\t\t\tlevel_start = level_end + 1;',
        '\t\t\tlevel_end += level_size;',
        '\t\t}',
        '\t}',
        '\tprintf("\\n");',
        '}',
    ]),
    ("트리 출력 - print_tree() 테스트", [
        "int main()",
        "{",
        "\tscreen_width = 50;",
        '\tchar data[10] = "ABCDEFG";',
        '\tprint_tree(data, 7, screen_width);',
        '\treturn 0;',
        "}",
    ]),
    ("heap sort - 최대힙 만들기 계획 - swap 함수", [
        "void swap(char *x, char *y)",
        "{",
        "\tchar t = *x;",
        "\t*x = *y;",
        "\t*y = t;",
        "}",
        "int main()",
        "{",
        '\tchar data[10] = "ABCDEFG";',
        '\tswap(data + 5, data + 6);',
        '\tprint_tree(data, 7, 32);',
        "}",
    ]),
    ("heap sort - 최대힙 만들기 계획 - parent 함수", [
        "int parent(int child)",
        "{",
        "\treturn (child - 1) / 2;",
        "}",
        "int main()",
        "{",
        '\tchar data[10] = "ABCDEFG";',
        '\tswap(data + parent(6), data + 6);',
        '\tswap(data + parent(3), data + 3);',
        '\tprint_tree(data, 7, 32);',
        "}",
    ]),
    ("heap sort - 최대힙 만들기 계획 - left_child 함수", [
        "int left_child(int parent)",
        "{",
        "\treturn 2 * parent + 1;",
        "}",
        "int main()",
        "{",
        '\tchar data[10] = "ABCDEFG";',
        '\tswap(data + left_child(2), data + 2);',
        '\tswap(data + left_child(1), data + 1);',
        '\tprint_tree(data, 7, 32);',
        "}",
    ]),
    ("heap sort - 최대힙 만들기 계획 - 최대 자식 찾기", [
        "int main()",
        "{",
        '\tchar data[10] = "ABCDEFG";',
        '\tint r = parent(6);',
        '\tint c = left_child(r);',
        '\tif (data[c] < data[c + 1])',
        '\t\tc = c + 1;',
        '\tif (data[r] < data[c])',
        '\t\tswap(data + r, data + c);',
        '\tprint_tree(data, 7, 32);',
        "}",
    ]),
    ("heap sort - 최대힙 만들기 계획 - level 2 이하까지", [
        "int main()",
        "{",
        '\tchar data[10] = "ABCDEFG";',
        '\tint r = 2;',
        '\tint c = left_child(r);',
        '\tif (data[c] < data[c + 1])',
        '\t\tc = c + 1;',
        '\tif (data[r] < data[c])',
        '\t\tswap(data + r, data + c);',
        '\tr = 1;',
        '\tc = left_child(r);',
        '\tif (data[c] < data[c + 1])',
        '\t\tc = c + 1;',
        '\tif (data[r] < data[c])',
        '\t\tswap(data + r, data + c);',
        '\tprint_tree(data, 7, 32);',
        "}",
    ]),
    ("heap sort - 최대힙 만들기 계획 - level2 이하 코드를 짧게", [
        "int main()",
        "{",
        '\tchar data[10] = "ABCDEFG";',
        '\tint i = 2, r, c;'
        '\tfor ( ; i >= 1; i -= 1) {',
        '\t\tr = i;',
        '\t\tc = left_child(r);',
        '\t\tif (data[c] < data[c + 1])',
        '\t\t\tc = c + 1;',
        '\t\tif (data[r] < data[c])',
        '\t\t\tswap(data + r, data + c);',
        '\t}'
        "}",
    ]),
    ("heap sort - 최대힙 만들기 계획 - level 1", [
        "int main()",
        "{",
        '\tchar data[10] = "ABCDEFG";',
        '\tint i = 2, r, c;'
        '\tfor ( ; i >= 1; i -= 1) {',
        '\t\tr = i;',
        '\t\tc = left_child(r);',
        '\t\tif (data[c] < data[c + 1])',
        '\t\t\tc = c + 1;',
        '\t\tif (data[r] < data[c])',
        '\t\t\tswap(data + r, data + c);',
        '\t}',
        '\tr = 0;',
        '\tc = left_child(r);',
        '\tif (data[c] < data[c + 1])',
        '\t\tc = c + 1;',
        '\tif (data[r] < data[c])',
        '\t\tswap(data + r, data + c);',
        '\tr = c;'
        '\tc = left_child(r);',
        '\tif (data[c] < data[c + 1])',
        '\t\tc = c + 1;',
        '\tif (data[r] < data[c])',
        '\t\tswap(data + r, data + c);',
        '\tprint_tree(data, 7, 32);',
        "}",
    ]),
    ("heap sort - 최대힙 만들기 계획 - level 1도 짧게", [
        "int main()",
        "{",
        '\tchar data[10] = "ABCDEFG";',
        '\tint i = 2, r, c;'
        '\tfor ( ; i >= 1; i -= 1) {',
        '\t\tfor (r = i; left_child(r) < 7; r = c) {', 
        '\t\t\tc = left_child(r);',
        '\t\t\tif (data[c] < data[c + 1])',
        '\t\t\t\tc = c + 1;',
        '\t\t\tif (data[r] >= data[c])',
        '\t\t\t\tbreak;',
        '\t\t\tswap(data + r, data + c);',
        '\t\t}',
        '\t}',
        '\tprint_tree(data, 7, 32);',
        "}",
    ]),
    ("heap sort - 최대힙 만들기 - right_child가 없을 때 처리", [
        "int main()",
        "{",
        '\tchar data[10] = "ABCDEFG";',
        '\tint i = 2, r, c;'
        '\tfor ( ; i >= 1; i -= 1) {',
        '\t\tfor (r = i; left_child(r) < 7; r = c) {', 
        '\t\t\tc = left_child(r);',
        '\t\t\tif (c + 1 < 7 && data[c] < data[c + 1])',
        '\t\t\t\tc = c + 1;',
        '\t\t\tswap(data + r, data + c);',
        '\t\t}',
        '\t}',
        '\tprint_tree(data, 7, 32);',
        "}",
    ]),
    ("heap sort - 최대힙 만들기", [
        "int main()",
        "{",
        '\tchar data[10] = "ABCDEFG";',
        '\tint n = 7;',
        "",
        "\tint i = parent(n - 1), c, r;",
        "\tfor ( ; i >= 0; i -= 1) {",
        "\t\tfor (r = i; r * 2 < n; r = c) {",
        "\t\t\tc = left_child(r);",
        "\t\t\tif (c < n - 1 && data[c] < data[c + 1])",
        "\t\t\t\tchild += 1;"
        "\t\t\tif (data[r] >= data[c])",
        "\t\t\t\tbreak;",
        "\t\t\tswap(data + r, data + c);",
        "\t\t}",
        "\t}",
        '\tprint_tree(data, 7, 32);',
        "}",
    ]),
    ("heap sort - 정렬하기 - G 처리", [
        "int main()",
        "{",
        '\tchar data[10] = "GEFDBAC";',
        '\tint n = 7;',
        "",
        '\tswap(data, data + 6);',
        "\tn = 6;",
        "\tint r = 0;",
        "\tint c = left_child(r);",
        "\tif (c < n - 1 && data[c] < data[c + 1])",
        "\t\tchild += 1;"
        "\tif (data[r] < data[c])",
        "\t\tswap(data + r, data + c);",
        "\tr = c;",
        "\tc = left_child(r);",
        "\tif (c < n - 1 && data[c] < data[c + 1])",
        "\t\tchild += 1;"
        "\tif (data[r] < data[c])",
        "\t\tswap(data + r, data + c);",
        '\tprint_tree(data, 7, 32);',
        "}",
    ]),

    ("heap sort - 정렬하기 - G 처리 코드를 짧게", [
        "int main()",
        "{",
        '\tchar data[10] = "GEFDBAC";',
        '\tint n = 7;',
        "",
        '\tswap(data, data + 6);',
        "\tn = 6;",
        "\tint c, r;",
        "\tfor (r = 0; left_child(r) < n; r = c) {"
        "\t\tc = left_child(r);",
        "\t\tif (c < n - 1 && data[c] < data[c + 1])",
        "\t\t\tchild += 1;"
        "\t\tif (data[r] >= data[c])",
        "\t\t\tbreak;",
        "\t\tswap(data + r, data + c);",
        "\t}",
        '\tprint_tree(data, 7, 32);',
        "}",
    ]),

    ("heap sort - 정렬하기 - F 처리", [
        "int main()",
        "{",
        '\tchar data[10] = "GEFDBAC";',
        '\tint n = 7;',
        "",
        '\tswap(data, data + 6);',
        "\tn = 6;",
        "\tint c, r;",
        "\tfor (r = 0; left_child(r) < n; r = c) {"
        "\t\tc = left_child(r);",
        "\t\tif (c < n - 1 && data[c] < data[c + 1])",
        "\t\t\tchild += 1;"
        "\t\tif (data[r] >= data[c])",
        "\t\t\tbreak;",
        "\t\tswap(data + r, data + c);",
        "\t}",
        "",
        '\tswap(data, data + 5);',
        "\tn = 5;",
        "\tint c, r;",
        "\tfor (r = 0; left_child(r) < n; r = c) {"
        "\t\tc = left_child(r);",
        "\t\tif (c < n - 1 && data[c] < data[c + 1])",
        "\t\t\tchild += 1;"
        "\t\tif (data[r] >= data[c])",
        "\t\t\tbreak;",
        "\t\tswap(data + r, data + c);",
        "\t}",
        '\tprint_tree(data, 7, 32);',
        "}",
    ]),

    ("heap sort - 정렬하기 - F 처리 코드를 짧게", [
        "int main()",
        "{",
        '\tchar data[10] = "GEFDBAC";',
        '\tint n = 7;',
        "",
        '\tint i, c, r;'
        '\tfor (i = n - 1; i >= 0; i -= 1) {',
        '\t\tswap(data, data + i);',
        "\t\tfor (r = 0; left_child(r) < i; r = c) {"
        "\t\t\tc = left_child(r);",
        "\t\t\tif (c < i - 1 && data[c] < data[c + 1])",
        "\t\t\t\tc = c + 1;"
        "\t\t\tif (data[r] >= data[c])",
        "\t\t\t\tbreak;",
        "\t\t\tswap(data + r, data + c);",
        '\t\t}',
        "\t}",
        "",
        '\tprint_tree(data, 7, 32);',
        "}",
    ]),
    
    ("heap sort - 정렬 함수", [
        "void sort(char *data, int n)",
        "{",
        "\tint i = parent(n - 1), c, r;",
        "\t/* heapify */",
        "\tfor ( ; i >= 0; i -= 1) {",
        "\t\tfor (r = i; left_child(r) < n; r = c) {",
        "\t\t\tc = left_child(r);",
        "\t\t\tif (c < n - 1 && data[c] < data[c + 1])",
        "\t\t\t\tc += 1;"
        "\t\t\tif (data[r] >= data[c])",
        "\t\t\t\tbreak;",
        "\t\t\tswap(data + r, data + c);",
        "",
        "\t/* sort */",
        "\tfor (i = n - 1; i >= 0; i -= 1) {",
        "\t\tswap(data, data + i);",
        "\t\tfor (r = 0; left_child(r) < i; r = c) {",
        "\t\t\tc = left_child(r);",
        "\t\t\tif (c < i - 1 && data[c] < data[c + 1])",
        "\t\t\t\tc += 1;",
        "\t\t\tif (data[r] >= data[c])",
        "\t\t\t\tbreak;",
        "\t\t\tswap(data + r, data + c);",
        "\t\t}",
        "\t}",
        "}",
        "int main()",
        "{",
        '\tchar data[10] = "GFEDCBA";',
        '\tint n = 7;',
        "",
        '\tsort(data, n);',
        "",
        '\tprint_tree(data, 7, 32);',
        "}",
    ]),
]

# 슬라이드 생성
for title, bullets in slides:
    add_slide(title, bullets)


output_path = "ppt/자료구조_7주차_14회차_PPT.pptx"
prs.save(output_path)
output_path
