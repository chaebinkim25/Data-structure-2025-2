from pptx import __version__ as pptx_version
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

print(pptx_version)

# 새 프레젠테이션 생성
prs = Presentation("template.pptx")
prs.slides.add_slide(prs.slide_masters[0].slide_layouts[0])  # 제목 슬라이드 추가
# 제목 슬라이드 설정
title_slide = prs.slides[0]
title_slide.shapes.title.text = "자료구조 (Data Structure\n9주차: 그래프 최단경로"


layout = prs.slide_masters[0].slide_layouts[1]  # 제목과 내용 레이아웃 선택

# 슬라이드 추가 함수
def add_slide(title, bullets):
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    content = slide.placeholders[1]
    content.text = bullets[0] if bullets else ""
    for bullet in bullets[1:]:
        p = content.text_frame.add_paragraph()
        p.text = bullet

slides = [
    ("그래프", [
        "그래프는 노드와 연결선을 모은 것이다","",
        "\t - vertexes: 전체 노드 리스트","",
        "\t - edges: 전체 연결선 리스트","",
        "\t - adj_matrix: 연결선 이차원 배열",
    ]),
    ("그래프 노드", [
        "그래프 노드는 노드의 정보와 이웃 리스트를 갖고 있다","",
        "\t - distance: 특정 노드로부터의 최단 거리","",
        "\t - parent: 최단 경로 상에서 이전 노드","",
        "\t - status: 방문 상태","",
        "\t - adj_list: 연결선 이차원 리스트",
    ]),
    ("그래프 연결선", [
        "그래프 연결선은 연결하는 노드와 연결선 정보를 갖고 있다","",
        "\t - length: 연결선의 비용","",
        "\t - node1: 연결선의 노드","",
        "\t - node2: 연결선의 노드","",
    ]),
    ("이번 시간 목차", [
        "그래프의 한 시작점에서 다른 모든 노드까지 최단거리","",
        "\t - 연결선 비용이 모두 1일 때","",
        "\t - 연결선 비용이 0 또는 1일 때","",
        "\t - 연결선 비용이 0 또는 양수일 때","",
        "\t - 연결선 비용이 0 또는 양수일 때 - 배열","",
        "\t - 연결선 비용에 음수가 있을 때","",
    ]),
    ("연결선 비용이 모두 1일 때", [
        "최단거리가 0인 노드는?","",
        "최단거리가 1인 노드는?","",
        "최단거리가 n인 노드는?","",
    ]),
    ("연결선 비용이 모두 1일 때", [
        "최단거리가 n인 노드들을 전부 알고 있을 때, 최단거리가 n + 1인 노드들을 찾는 방법은?","",
    ]),
    ("연결선 비용이 모두 1일 때", [
        "함수 BFS(그래프 G, 시작 노드 s):",
        "to_visit = 새로운 큐",
        "반복: G의 모든 노드 u에 대해,",
        "\tu->distance = 아주 큰 수",
        "\tu->parent = 없음",
        "\tu->status = 추가된적없음",
        "s->distance = 0",
        "s->status = 추가된적있음",
        "to_visit에 s 추가",
        "반복: to_visit에 데이터가 있는 동안",
        "\tu = DEQUEUE(to_visit)",
        "\t반복: u의 모든 이웃 노드 v에 대해,",
        "\t\t만약: v->status가 추가된적없음이면,",
        "\t\t\tv->parent = u",
        "\t\t\tv->distance = u->distance + 1",
        "\t\t\tv->status = 추가된적있음",
        "\t\t\tto_visit에 v 추가",
    ]),
    ("연결선 비용이 모두 1일 때", [
        "다음 그래프에서 BFS(G, s)를 실행하면?",
    ]),
    ("연결선 비용이 0 또는 1일 때", [
        "최단거리가 0인 노드는?","",
        "최단거리가 1인 노드는?","",
        "최단거리가 n인 노드는?","",
    ]),
    ("연결선 비용이 0 또는 1일 때", [
        "최단거리를 구한 노드들로부터 다른 노드들의 최단 거리를 구하는 방법은?","",
    ]),
    ("연결선 비용이 0 또는 1일 때", [
        "양방향 큐","",
        "\t - PUSH_FRONT(): 제일 앞에 데이터 추가","",
        "\t - PUSH_REAR(): 제일 뒤에 데이터 추가","",
        "\t - POP_FRONT: 제일 앞 데이터 삭제","",
        "\t - POP_REAR: 제일 뒤 데이터 삭제","",
    ]),
    ("연결선 비용이 0 또는 1일 때", [
        "함수 BFS_0_1(그래프 G, 시작 노드 s):",
        "to_visit = 새로운 양방향 큐",
        "반복: G의 모든 노드 u에 대해,",
        "\tu->distance = 아주 큰 수",
        "\tu->parent = 없음",
        "\tu->status = 처리전",
        "s->distance = 0",
        "PUSH_FRONT(to_visit, s)",
        "반복: to_visit에 데이터가 있는 동안",
        "\tu = POP_FRONT(to_visit)",
        "\t만약: u->status == 처리됨이면,",
        "\t\t다음 반복으로",
        "\tu->status = 처리됨",
        "\t반복: u의 모든 이웃 노드 v와 (u, v)를 연결하는 e에 대해,",
        "\t\t만약: v->distance ≤ u->distance + e->length면,",
        "\t\t\t다음 반복으로",
        "\t\tv->parent = u",
        "\t\tv->distance = u->distance + e->length",
        "\t\t만약: e->length == 0이면,"
        "\t\t\tPUSH_FRONT(to_visit, v)",
        "\t\t그 외에:"
        "\t\t\tPUSH_REAR(to_visit, v)",
    ]),
    ("연결선 비용이 0 또는 1일 때", [
        "다음 그래프에서 BFS_0_1(G, s)를 실행하면?",
    ]),
    ("연결선 비용이 0 또는 양수일 때", [
        "최단거리가 가장 작은 노드는?","",
        "최단거리가 두번째로 작은 노드는?","",
        "최단거리가 세번째로 작은 노드는?","",
    ]),
    ("연결선 비용이 0 또는 양수일 때", [
        "최단거리를 구한 노드들로부터 다른 노드들의 최단 거리를 구하는 방법은?","",
    ]),
    ("연결선 비용이 0 또는 양수일 때", [
        "함수 DIJKSTRA(그래프 G, 시작 노드 s):",
        "to_visit = 우선순위 큐 (최소 length)",
        "반복: G의 모든 노드 u에 대해,"
        "\tu->distance = 아주 큰 수"
        "\tu->parent = 없음"
        "\tu->status = 방문전",
        "s->distance = 0",
        "s->status = 방문예정",
        "to_visit에 s 추가",
        "반복: to_visit에 데이터가 있는 동안"
        "\tu = DEQUEUE(to_visit)"
        "\t만약: u->status != 방문예정이면,"
        "\t\t다음 반복으로",
        "\tu->status = 방문중",
        "\t반복: u의 모든 이웃 노드 v와 (u, v)를 연결하는 e에 대해,",
        "\t\t만약: v->distance ≤ u->distance + e->length이면,",
        "\t\t\t다음 반복으로"
        "\t\tv->parent = u",
        "\t\tv->distance = u->distance + e->length",
        "\t\tv->status = 방문예정",
        "\t\tENQUEUE(to_visit, v)",
        "\tu->status = 방문후",
    ]),
    ("연결선 비용이 0 또는 양수일 때", [
        "다음 그래프에서 DIJKSTRA(G, s)를 실행하면?",
    ]),
    ("연결선 비용이 0 또는 양수일 때", [
        "함수 DIJKSTRA_ARRAY(그래프 G, 시작 노드 s):",
        "반복: G의 모든 노드 u에 대해,"
        "\tu->distance = 아주 큰 수"
        "\tu->parent = 없음"
        "\tu->status = 방문전",
        "s->distance = 0",
        "반복: G에 status가 방문전인 노드가 있는 동안"
        "\tu = G의 status가 방문전인 노드들 중 distance가 최소"
        "\tu->status = 방문후",
        "\t반복: u의 모든 이웃 노드 v와 (u, v)를 연결하는 e에 대해,",
        "\t\t만약: v->distance ≤ u->distance + e->length이면,",
        "\t\t\t다음 반복으로",
        "\t\tv->parent = u",
        "\t\tv->distance = u->distance + e->length",
    ]),
    ("연결선 비용이 0 또는 양수일 때", [
        "다음 그래프에서 DIJKSTRA_ARRAY(G, s)를 실행하면?",
    ]),
    ("일반적인 경우", [
        "1회차에서 시작점부터 노드 A까지 최단거리는?","",
        "2회차에서 시작점부터 노드 A까지 최단거리는?","",
        "3회차에서 시작점부터 노드 A까지 최단거리는?","",
    ]),
    ("일반적인 경우", [
        "최대 반복해야할 횟수는?","",
    ]),
    ("일반적인 경우", [
        "함수 BELLMAN_FORD(그래프 G, 시작 노드 s):",
        "반복: G의 모든 노드 u에 대해,",
        "\tu->distance = 아주 큰 수",
        "\tu->parent = 없음",
        "s->distance = 0",
        "반복: G의 노드 개수 - 1번",
        "\t반복: G의 모든 연결선 e에 대해,",
        "\t\tu = e->node1",
        "\t\tv = e->node2"
        "\t\t만약: u->distance < 아주 큰수이고, v->distance > u->distance + e->length이면,",
        "\t\t\tv->distance = u->distance + e->length",
        "\t\t\tv->parent = u",
        "\t\t만약: v->distance < 아주 큰수이고, u->distance > v->distance + e->length이면,",
        "\t\t\tu->distance = v->distance + e->length",
        "\t\t\tu->parent = v",
    ]),
    ("연결선 비용이 0 또는 양수일 때", [
        "다음 그래프에서 BELLMAN_FORD(G, s)를 실행하면?",
    ]),
    ("일반적인 경우", [
        "음수 사이클이 있다면?","",
    ]),
    ("일반적인 경우", [
        "음수 사이클이 있는 경우와 없는 경우 |V|번째 반복에서 일어나는 일은?","",
    ]),
    ("일반적인 경우", [
        "음수 사이클을 찾는 방법은?","",
    ]),
    ("일반적인 경우", [
        "함수 BELLMAN_FORD_SAFE(그래프 G, 시작 노드 s):",
        "반복: G의 모든 노드 u에 대해,",
        "\tu->distance = 아주 큰 수",
        "\tu->parent = 없음",
        "\tu->status = 깨끗함",
        "s->distance = 0",
        "반복: G의 노드 개수 - 1번",
        "\t반복: G의 모든 연결선 e에 대해,",
        "\t\tu = e->node1",
        "\t\tv = e->node2"
        "\t\t만약: u->distance < 아주 큰수이고, v->distance > u->distance + e->length이면,",
        "\t\t\tv->distance = u->distance + e->length",
        "\t\t\tv->parent = u",
        "\t\t만약: v->distance < 아주 큰수이고, u->distance > v->distance + e->length이면,",
        "\t\t\tu->distance = v->distance + e->length",
        "\t\t\tu->parent = v",

        "bad_roots = 새로운 리스트"
        "반복: G의 모든 연결선 e에 대해,",
        "\tu = e->node1",
        "\tv = e->node2"
        "\t\t만약: u->distance < 아주 큰수이고, v->distance > u->distance + e->length이면,",
        "\t\t\tbad_roots에 v가 없으면 v 추가",
        "\t\t만약: v->distance < 아주 큰수이고, u->distance > v->distance + e->length이면,",
        "\t\t\tbad_roots에 u가 없으면 u 추가",
        "bad_roots의 모든 노드에 대해 DFS로 방문된 노드의 status를 오염됨으로 변경"
    ]),


    ("요약", [
        "너비 우선 탐색: 연결선의 비용이 모두 1일 때, 최단 거리 계산","",
        "DIJKSTRA: 연결선의 비용이 0 또는 양수일 때",
        "BELLMAN-FORD: 연결선의 비용에 음수가 있을 때도",
    ]),


]

# 슬라이드 생성
for title, bullets in slides:
    add_slide(title, bullets)


def add_table_slides(title, bullets, table_data):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title

        content = slide.placeholders[1]
        content.text = bullets[0] if bullets else ""
        for bullet in bullets[1:]:
                p = content.text_frame.add_paragraph()
                p.text = bullet

        cols, rows = len(table_data), len(table_data[0]) if table_data else 0
        left = Inches(0.5)
        top = Inches(5)
        width = Inches(9)
        height = Inches(0.8 + 0.3 * rows)
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        for c in range(cols):
                for r in range(rows):
                        table.cell(r, c).text = table_data[c][r]


# 저장
output_path = "ppt/자료구조_9주차_17회차_PPT.pptx"
prs.save(output_path)
output_path
