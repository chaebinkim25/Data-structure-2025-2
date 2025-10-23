from pptx import __version__ as pptx_version
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

print(pptx_version)

# 새 프레젠테이션 생성
prs = Presentation("template.pptx")
prs.slides.add_slide(prs.slide_masters[0].slide_layouts[0])  # 제목 슬라이드 추가
# 제목 슬라이드 설정
title_slide = prs.slides[0]
title_slide.shapes.title.text = "자료구조 (Data Structure)\n5주차: 트리"


layout = prs.slide_masters[0].slide_layouts[1]  # 제목과 내용 레이아웃 선택

# 슬라이드 추가 함수
def add_slide(title, bullets):
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    content = slide.placeholders[1]
    content.text = bullets[0] if bullets else ""
    for bullet in bullets[1:]:
        p = content.text_frame.add_paragraph()
        p.text = bullet

slides = [

    ("지난 시간 요약", [
        "트리: 루트로부터 부모 자식 관계로 연결된 구조","",
        "이진 트리: 자식 노드가 최대 2개인 트리","",
        "이진 트리 예제: 수식 트리","",
        "이진 트리 노드: left, right, value","",
        "트리 방문: Preorder, Inorder, Postorder","",
    ]),


    ("오늘 목표", [
        "1) 이진 트리 및 방문 구현","",
        "2) 예제: 방문 순서로 이진 트리 재구성","",
        "3) 예제: 수식 트리","",
    ]),

    ("이진 트리 및 방문 구현", [
        "이진 트리 복습","",
        "이진 트리 노드 타입",
        "이진 트리 타입","",
        "이진 트리 방문",
        "이진 트리 삭제하기",
    ]),

    ("이진 트리", [
        "자식 노드가 최대 2개인 트리","",
        "        - root: 트리의 시작점 (타입: 이진 트리 노드)","",
    ]),

    ("이진 트리 노드 타입", [
        "데이터와 자식 노드 2개 정보의 묶음","",
        "        - left: 왼쪽 자식 노드 주소","",
        "        - right: 오른쪽 자식 노드 주소","",
        "        - value: 저장된 데이터","",
    ]),

    ("이진 트리 노드 타입", [
        "struct bi_node {",
        "        struct bi_node *left;",
        "        struct bi_node *right;",
        "        char value;",
        "};",
    ]),

    ("이진 트리 노드 생성 함수", [
        "struct bi_node *new_bi_node(char v, struct bi_node *l, struct bi_node *r) "
        "{",
        "       struct bi_node *node = (struct bi_node *)malloc(sizeof(struct bi_node));",
        "       node->left = l;",
        "       node->right = r;",
        "       node->value = v;",
        "       return node;",
        "}",
    ]),

    ("이진 트리", [
        "struct bi_root {",
        "        struct bi_node *root;",
        "};",
    ]),

    ("이진 트리 생성 함수", [
        "struct bi_root *new_bi_root(struct bi_node *root)",
        "{",
        "       struct bi_root *tree = (struct bi_root *)malloc(sizeof(struct bi_root));",
        "       tree->root = root;",
        "       return tree;",
        "}",
    ]),


    ("이진 트리 예시", [
        "다음과 같은 이진 트리를 만드는 함수를 써보자",
        "       A",
        "      / \\",
        "     B   C",
        "    / \\",
        "   D   E",
        "struct bi_root *sample_tree() ",
        "{",
        "       struct bi_node *D = new_bi_node('D', NULL, NULL);",
        "       struct bi_node *E = new_bi_node('E', NULL, NULL);",
        "       struct bi_node *C = new_bi_node('C', NULL, NULL);",
        "       struct bi_node *B = new_bi_node('B', D, E);",
        "       struct bi_node *A = new_bi_node('A', B, C);",
        "       struct bi_root *T = new_bi_root(A);",
        "       return T;",
        "}",
    ]),

    ("깊이 우선 이진 트리 방문", [
        "함수 TRAVERSE(node):","",
        "만약: node != NULL이면",
        "       pre_task(node->value)",
        "       TRAVERSE(node->left)",
        "       in_task(node->value)",
        "       TRAVERSE(node->right)",
        "       post_task(node->value)",
    ]),

    ("깊이 우선 이진 트리 방문", [
        "void traverse(struct bi_node *node) ",
        "{",
        "       if (node != NULL) {",
        "               pre_task(node->value);",
        "               traverse(node->left);",
        "               in_task(node->value);",
        "               traverse(node->right);",
        "               post_task(node->value);",
        "       }",
        "}",
    ]),

    ("이진 트리 방문시 작업 타이밍", [
        "다음 트리에 대한 traverse(root)에서",
        "pre-, in-, post-order 작업이 실행되는 순서는?",
        "",
        "       A",
        "      / \\",
        "     B   C",
        "    / \\",
        "   D   E",
        "",
        "preorder: A B D E C",
        "inorder: D B E A C",
        "postorder: D E B C A",
    ]),

    ("이진 트리 삭제하기", [
        "트리 방문시 노드를 free할 자연스러운 작업 타이밍은?","",
        "(힌트: 먼저 변수를 다 쓰고 free한다)",
    ]),

    ("이진 트리 삭제하기", [
        "함수 DELETE(node):","",
        "만약: node != NULL이면",
        "       DELETE(node->left)",
        "       DELETE(node->right)",
        "       free(node)",
    ]),

    ("이진 트리 삭제하기", [
        "void bi_node_delete(struct bi_node *node) ", 
        "{",
        "       if (node != NULL) {",
        "               bi_node_delete(node->left);",       
        "               bi_node_delete(node->right);",
        "               free(node);",
        "       }",
        "}",
    ]),

    ("이진 트리 삭제하기", [
        "void bi_tree_delete(struct bi_root *tree) ",
        "{",
        "       bi_node_delete(tree->bi_node);",
        "       free(tree);",
        "}",
    ]),

    ("이진 트리 및 방문 구현", [
        "이진 트리 - 자식 노드가 2개까지 가능한 트리","",
        "이진 트리 노드 타입 - left, right, value",
        "이진 트리 타입 - root","",
        "이진 트리 방문 - 순서 preorder, inorder, postorder",
        "이진 트리 삭제하기 - postorder로",
    ]),

    ("예제: 방문 순서로 이진 트리 재구성", [
        "트리 재구성 퀴즈","",
        "트리 재구성 전략",
        "\t - 부분 리스트",
        "\t - 데이터 찾기",
        "\t - 부분 트리 찾기","",
        "트리 재구성 C코드",
    ]),

    ("트리 재구성 퀴즈", [
        "preorder: A B D E C",
        "inorder: D B E A C",
        "Original tree???",
        "",
        "       A",
        "      / \\",
        "     B   C",
        "    / \\",
        "   D   E",
        "",
    ]),

    ("트리 재구성 전략", [
        "함수 rebuild(preorder, inorder):",
        "만약: preorder가 비어있으면, 반환: NULL",
        "root_val = RETRIEVE(preorder, 1)",
        "in_pivot = inorder에서 root_val의 순서",
        "in_l = inorder 부분리스트(inorder 시작, in_pivot - 1)",
        "in_r = inorder 부분리스트(in_pivot + 1, inorder 끝)",
        "pre_l = preorder에서 in_l에 해당하는 부분리스트",
        "pre_r = preorder에서 in_r에 해당하는 부분리스트",
        "left = rebuild(pre_l, in_l)",
        "right = rebuild(pre_r, in_r)",
        "반환: 새 노드(root_val, left, right)",
    ]),

    ("트리 재구성 전략", [
        "preorder: A B D E C",
        "inorder: D B E A C",
        "rebuild(preorder, inorder) 실행 과정은?",
    ]),

    ("트리 재구성 전략 - 부분 리스트", [
        "preorder, inorder를 저장하기에 적합한 자료구조는?",
    ]),

    ("트리 재구성 전략 - 부분 리스트", [
        "struct view {",
        "       const char *data;",
        "       int start;",
        "       int end;",
        "       int size;",
        "};",
    ]),

    ("트리 재구성 전략 - 부분 리스트", [
        "struct view *new_view(const char *arr, int from, int to)",
        "{",
        "       struct view *obj = (struct view *)malloc(sizeof(struct view));",
        "       obj->data = arr;",
        "       obj->start = from;",
        "       obj->end = to;",
        "       obj->size = to - from + 1;",
        "       return obj;",
        "}",
    ]),

    ("트리 재구성 전략 - 부분 리스트", [
        "preorder: A B D E C",
        "preorder의 3 - 5번에 해당하는 부분 리스트를 new_view 함수로 만들려면?",
    ]),

    ("트리 재구성 전략 - 데이터 찾기", [
        "preorder: A B D E C",
        "preorder에서 B의 순서를 찾는 방법은?",
    ]),

    ("이진 트리 재구성", [
        "int find_index(struct view *range, char value) ",
        "{",
        "       for (int i = range->start; i <= range->end; i++) {",
        "               if (range->data[i] == value) return i;",
        "       }",
        "       return -1;",
        "}",
    ]),

    ("트리 재구성 전략 - 데이터 찾기", [
        "preorder: A B D E C",
        "find_index(preorder, 'D')의 실행 과정은?",
    ]),

    ("트리 재구성 전략 - 부분 트리 찾기", [
        "inorder: D B E A C",
        "inorder에서 루트의 왼쪽에 해당하는 부분리스트를 구하면?",
        "",
        "       A",
        "      / \\",
        "     B   C",
        "    / \\",
        "   D   E",
        "",
    ]),

    ("트리 재구성 전략 - 부분 트리 찾기", [
        "preorder: A B D E C",
        "preorder에서 부분 트리 B에 해당하는 부분리스트를 구하면?",
        "",
        "       A",
        "      / \\",
        "     B   C",
        "    / \\",
        "   D   E",
        "",
    ]),

    ("트리 재구성 전략 - 부분 트리 찾기", [
        "preorder: A B D E C",
        "preorder에서 D B E로 이뤄진 부분리스트를 구하면? (단, D B E는 부분 트리를 이룬다)",
    ]),

    ("트리 재구성 전략 - 부분 트리 찾기", [
        "struct view *find_subrange(struct view *range, struct view *set) ",
        "{",
        "       int is_in_set[256] = {0};",
        "       for (int i = set->start; i <= set->end; i++) {",
        "               int key = (unsigned)set->data[i];",
        "               is_in_set[key] = 1;",
        "       }",

        "       int i = range->start;",
        "       while (i <= range->end) {",
        "               int key = (unsigned)range->data[i];",
        "               if (is_in_set[key])",
        "                       break;",
        "               i++;",
        "       }",

        "       return new_view(range->data, i, i + set->size - 1);",
        "}",
    ]),

    ("트리 재구성 전략 - 부분 트리 찾기", [
        "preorder: A B D E C",
        "subrange_set: D B E",
        "find_subrange(preorder, subrange_set)의 실행 과정은?",
    ]),

    ("트리 재구성 C 코드", [
        "rebuild(preorder, inorder) 함수를 C로 구현하면?",
    ]),

    ("이진 트리 재구성", [
        "struct bi_node *rebuild(struct view *pre, struct view *in) ",
        "{",
        "       if (pre->size == 0 || in->size == 0) return NULL;",
        "       int root_val = pre->data[pre->start];",
        "       int in_pivot = find_index(in, root_val);",
        "       struct view *in_l = new_view(in->data, in->start, in_pivot - 1);",
        "       struct view *in_r = new_view(in->data, in_pivot + 1, in->end);",
        "       struct view *pre_l = find_range(pre->data, pre->start, pre->end, in_l);",
        "       struct view *pre_r = find_range(pre->data, pre->start, pre->end, in_r);",
        "       struct bi_node *left = rebuild(pre_l, in_l);",
        "       struct bi_node *right = rebuild(pre_r, in_r);",
        "       free(pre_l); free(pre_r); free(in_l); free(in_r);",
        "       return create_bi_node(root_val, left, right);",
        "}",
    ]),

    ("트리 재구성 C 코드", [
        "preorder: A B D E C",
        "inorder: D B E A C",
        "rebuild(preorder, inorder) 함수의 실행 과정은?",
    ]),

    ("예제: 방문 순서로 이진 트리 재구성 요약", [
        "트리 재구성 퀴즈","",
        "트리 재구성 전략 - 재귀적으로 왼쪽/오른쪽 구분",
        "\t - 부분 리스트 - 리스트의 일부 구간",
        "\t - 데이터 찾기",
        "\t - 부분 트리 찾기","",
        "트리 재구성 C코드",
    ]),

    ("다음 시간 예고", [
        "이진 트리의 특수한 케이스","",
    ]),

    ("예제: 수식 트리", [
        "수식 트리 퀴즈","",
        "수식 트리 생성",
        "\t - 노드 타입",
        "\t - C코드","",
        "수식 트리 계산",
    ]),

    ("수식 트리 퀴즈", [
        "1 + 2 * 3을 수식 트리로 만들고, 값을 계산하면?","",
    ]),

    ("수식 트리 생성", [
        "함수 expr_tree(수식 노드 큐 Q, 처리중인 우선순위 p):",
        "node = DEQUEUE(Q)",
        "반복: size(Q) > 0인 동안",
        "       op = FRONT(Q)",
        "       만약: 우선순위(op) ≤ p이면 반복 종료",
        "       DEQUEUE(Q)",
        "       op->left = node",
        "       op->right = expr_tree(Q, 우선순위(op))",
        "       node = op",
        "반환: node",
    ]),

    ("수식 트리 생성 - 노드 타입", [
        "수식 트리 노드의 데이터로 필요한 것은?","",
    ]),
    
    ("수식 트리 생성 - 노드 타입", [
        "숫자 또는 연산자가 저장된다","",
        "        - type: 숫자, 더하기, 곱하기 구분",
        "        - value: 숫자일 경우 값을 저장",
        "        - precedence: 연산자 우선순위",
        "                      더하기는 1",
        "                      곱하기는 2",
        "                      숫자는 -1 (연산자 아님)",
    ]),

    ("수식 트리 생성 - 노드 타입", [
        "struct expr_node {",
        "        struct expr_node *l;    // 왼쪽 자식 노드",
        "        struct expr_node *r;   // 오른쪽 자식 노드",
        "        char type;          // 'i': 숫자, '+': 더하기, '*': 곱하기",
        "        int value;          // 결과값",
        "        int precedence;     // 연산자 우선순위",
        "};"
    ]),

    ("수식 트리 생성 - C 코드", [
        "수식 노드 큐를 수식 트리로 만드는 함수를 C 코드로 써보면","",
    ]),

    ("예제: 더하기/곱하기 수식 트리 생성", [
        "struct expr_node *expr_tree(struct token_queue *q, int cur_prec)",
        "{",
        "       struct expr_node *node = expr_q_dequeue(q);",
        "       while (!expr_q_empty(q)) {",
        "               struct expr_node *op = expr_q_front(q);",
        "               if (op->precedence <= cur_prec)",
        "                       break;",
        "               expr_q_dequeue(q);",
        "               op->l = node;",
        "               op->r = expr_tree(q, op->precedence);",
        "               node = op;",
        "       }",
        "       return node;",
        "}",
    ]),

    ("예제: 더하기/곱하기 수식 트리 생성", [
        "expr_q: [1, '+', 2, '*', 3, '*', 4, '+', 5]",
        "expr_tree(expr_q, 0)의 실행 과정은?","",
        "       1. node = 1",
        "       2. op = '+'",
        "       3. 우선순위(op) > 0이므로 진행",
        "       4. DEQUEUE(Q)",
        "       5. op->left = node;",
        "       6. op->right = expr_tree(q, 1);",
        "               6-1. node = 2",
        "               6-2. op = '*'",
        "               6-3. 우선순위(op) > 1이므로 진행",
        "               6-4. DEQUEUE(Q)",
        "               6-5. op->left = node;",
        "               6-6. op->right = expr_tree(q, 2);",
        "                       6-6-1. node = 3",
        "                       6-6-2. op = '*'",
        "                       6-6-3. 우선순위(op) <= 2이므로 반복 종료",
        "                       6-6-4. 반환: node",
        "               6-7. op->right = 3;",
        "               6-8. node = op;",
        "               6-9. 다음 반복",
        "               6-10. op = '*'",
        "               6-11. 우선순위(op) > 1이므로 진행",
        "               6-12. DEQUEUE(Q)",
        "               6-13. op->left = node;",
        "               6-14. op->right = expr_tree(q, 2);",
        "                       6-14-1. node = 4",
        "                       6-14-2. op = '+'",
        "                       6-14-3. 우선순위(op) <= 2이므로 반복 종료",
        "                       6-14-4. 반환: node",
        "               6-15. op->right = 4;",
        "               6-16. node = op;",
        "               6-17. 다음 반복",
        "               6-18. op = '+'",
        "               6-19. 우선순위(op) <= 1이므로 반복 종료",
        "               6-20. 반환: node",
        "       7. op->right = node;",
        "       8. node = op;",
        "       9. 다음 반복",
        "       10. op = '+'",
        "       11. 우선순위(op) > 0이므로 진행",
        "       12. DEQUEUE(Q)",
        "       13. op->left = node;",
        "       14. op->right = expr_tree(q, 1);",
        "               14-1. node = 5",
        "               14-2. 반복 종료",
        "               14-3. 반환: node",
        "       15. op->right = 5;",
        "       16. node = op;",
        "       17. 다음 반복",
        "       18. 반복 종료",
        "       19. 반환: node",
        "}",
    ]),
    ("수식 트리 계산", [
        "수식 트리를 계산할 때 pre-, in-, post-order 중 적절한 작업 타이밍은?",
    ]),

    ("수식 트리 계산", [
        "함수 post_task(struct expr_node *expr)에서 할 일로 적절한 것은?","",
    ]),

    ("수식 트리 계산", [
        "함수 post_task(struct expr_node *expr):",
        "스위치(expr->type):",
        "케이스 '+':",
        "       expr->value = expr->l->value + expr->r->value;",
        "       휴식;",
        "케이스 '*':",
        "       expr->value = expr->l->value * expr->r->value;",
        "       휴식;",
        "케이스 'i':",
        "       // 숫자 노드는 이미 값이 들어있음",
        "       휴식;",
    ]),

    ("수식 트리 계산", [
        "void expr_post_task(struct expr_node *e) ",
        "{",
        "       if (e->type == '+') {",
        "               e->value = e->l->value + e->r->value;",
        "       } else if (e->type == '*') {",
        "               e->value = e->l->value * e->r->value;",
        "       }",
        "}",
    ]),

    ("수식 트리 계산", [
        "void expr_eval(struct expr_node *e)",
        "{",
        "        if (e->l)",
        "                expr_eval(e->l);",
        "        if (e->r)",
        "                expr_eval(e->r);",
        "        expr_post_task(e);",
        "}",
    ]),

    ("수식 트리 계산", [
        "1 + 2 * 3",
        "다음 수식 트리의 루트에서 expr_eval의 실행 과정은?","",
        "       1. left_val = expr_eval(1);",
        "               1-1. 반환: 1",
        "       2. right_val = expr_eval('*');",
        "               2-1. left_val = expr_eval(2);",
        "                       2-1-1. 반환: 2",
        "               2-2. right_val = expr_eval(3);",
        "                       2-2-1. 반환: 3",
        "               2-3. 반환: 2 * 3 = 6",
        "       3. 반환: 1 + 6 = 7",
    ]),
    



    ("C 코딩 스타일", [
        "리눅스 커널 코딩 스타일","",
        "이유: 세계에서 가장 큰 C 프로젝트","",
        "C 코딩을 매우 많이 하고, 많이 읽은 프로그래머들이 정한 스타일","",
        "https://github.com/torvalds/linux/blob/master/Documentation/process/coding-style.rst","",
    ]),
    ("C 코딩 스타일 - 들여쓰기", [
        "들여쓰기: 8칸","",
        "이유: 들여쓰기 되어 있음을 확실하게 볼 수 있다.","",
        '#include <stdio.h>\n\nint main()\n{\n        printf("Hello, world!");\n}'
    ]),
    ("C 코딩 스타일 - switch문의 들여쓰기", [
        "case 라벨의 들여쓰기: 없음","",
        "이유: 지나친 들여쓰기를 방지한다.","",
        'switch (suffix) {\n'
        "case 'G':\n"
        "case 'g':",
        "        mem <<= 30;\n"
        "        break;\n"
        "case 'M':\n"
        "case 'm':",
        "        mem <<= 20;\n"
        "        break;\n"
        "case 'K':\n"
        "case 'k':",
        "        mem <<= 10;\n"
        "        /* fallthrough */\n"
        "default:\n"
        "        break;\n"
        "}"
    ]),
    ("C 코딩 스타일 - 함수의 중괄호", [
        "중괄호 여는 위치: 헤더 다음 줄","",
        "이유: 함수는 특별하다. C에서는 함수 정의 안에 함수 정의를 못한다.","",
        "C 언어 제작자들이 쓰는 방식","",
        'int function(int x)\n'
        '{\n'
        '        /* body of function */\n'
        '}\n'
    ]),
    ("C 코딩 스타일 - 모든 문장 블록의 중괄호", [
        "중괄호 여는 위치: 줄의 오른쪽 끝","",
        "중괄호 닫는 위치: 줄의 왼쪽 끝 (대부분의 경우 다른 문자 없음)","",
        'if (x_is_true) {\n'
        '        do_y();\n'
        '        do_z();\n'
        '}\n'
    ]),
    ("C 코딩 스타일 - do-문의 중괄호", [
        "중괄호를 닫고 나서 같은 줄에 while을 쓴다.","",
        'do {\n        body of do-loop\n} while (condition);'
    ]),
    ("C 코딩 스타일 - if-문의 중괄호", [
        "중괄호를 닫고 나서 같은 줄에 else를 쓴다.","",
        'if (x == y) {\n        ..\n} else if (x > y) {\n        ...\n} else {\n        ....\n}'
    ]),
    ("C 코딩 스타일 - if-문의 중괄호", [
        "if와 else에서 모두 문장 하나만 실행할 때는 중괄호를 쓰지 않는다.","",
        'if (condition)\n        do_this();\nelse\n        do_that();'
    ]),
    ("C 코딩 스타일 - if-문의 중괄호 (나쁜 예)", [
        "if의 문장을 숨길 것이 아닌 한 한줄에 쓰지 않는다.","",
        'if (condition) do_this;\ndo_something_everytime;'
    ]),
    ("C 코딩 스타일 - if-문의 중괄호 (나쁜 예)", [
        "괄호를 쓰지 않기 위해서 쉼표를 쓰지 않는다.","",
        'if (condition)\n        do_this(), do_that();'
    ]),
    ("C 코딩 스타일 - if-문의 중괄호", [
        "if나 else 중 하나에만 문장이 두개 이상 있어도 모두 중괄호를 쓴다.","",
        'if (condition) {\n        do_this();\n        do_that();\n} else {\n        otherwise();\n}'
    ]),
    ("C 코딩 스타일 - 반복문의 중괄호", [
        "반복문은 단순 문장 한개만 있지 않은 한 항상 중괄호를 쓴다.","",
        'while (condition) {\n        if (test)\n                do_something();\n}'
    ]),
    ("C 코딩 스타일 - 띄어쓰기 - 키워드", [
        "다음 키워드 다음에는 띄어쓰기를 한칸 한다.","",
        "if, switch, case, for, do, while","",
        "sizeof 다음에는 띄어쓰기를 하지 않는다.","",
        "좋은 예: s = sizeof(struct file);","",
        "괄호 안쪽 주변에는 띄어쓰기를 하지 않는다.","",
        "나쁜 예: s = sizeof( struct file );","",
    ]),
    ("C 코딩 스타일 - 띄어쓰기 - 포인터", [
        "포인터 변수를 선언하거나, 포인터를 반환하는 함수를 선언할 때,","",
        "* 은 변수 이름이나 함수 이름에 붙인다.","",
        "* 을 타입 옆에 붙이지 않는다.","",
        "char *linux_banner;\nunsigned long long memparse(char *ptr, char **retptr);\nchar *match_strdup(substring_t *s);","",
    ]),
    ("C 코딩 스타일 - 띄어쓰기 - 이항연산자", [
        "다음 이항 연산자 / 삼항 연산자 주변에는 띄어쓰기를 한칸 한다.","",
        "= + - * / % < > <= >= == != & | ? :","",
        "다음 구조체 멤버 연산자 주변에는 띄어쓰기를 하지 않는다.","",
        ". ->","",
        "area = rect.x * rect.y;","",
    ]),
    ("C 코딩 스타일 - 띄어쓰기 - 단항연산자", [
        "다음 단항 연산자 뒤에는 띄어쓰기를 하지 않는다.","",
        "& * + - ~ ! sizeof","",
        "다음 postfix 증감 단항 연산자 주변에는 띄어쓰기를 하지 않는다.","",
        "++ --","",
        "다음 prefix 증감 단항 연산자 주변에는 띄어쓰기를 하지 않는다.","",
        "++ --","",
        "i++;","",
    ]),
    ("C 코딩 스타일 - 이름", [
        "C 언어는 단순하고 미니멀한 것을 추구한다.","",
        "어떤 프로그래밍 언어에서는 ThisVariableIsATemporaryCounter로 쓰지만","",
        "C 언어에서는 쓰기 훨씬 쉽고, 이해하기는 많이 어렵지 않은 tmp로 쓴다.","",
        "한편, 전역 범위에 있는 이름은 자세하게 만들어야 한다.","",
    ]),
    ("C 코딩 스타일 - 전역 변수 이름, 함수 이름", [
        "전역 변수는 다른 방법이 전혀 없어서 어쩔 수 없이 꼭 필요할 때만 쓴다.","",
        "전역 변수와 함수의 이름은 단어들을 _로 구분해서 자세하게 만든다.","",
        "좋은 예: count_active_users()","",
        "나쁜 예: cntusr()","",
    ]),
    ("C 코딩 스타일 - 지역 변수 이름", [
        "지역 변수 이름은 짧고, 정보를 담고 있어야 한다.","",
        "좋은 예: 루프 변수 i","",
        "나쁜 예: 루프 변수 loop_counter (오해의 소지가 없을 때)","",
        "좋은 예: 임시 변수 tmp","",
    ]),
    ("C 코딩 스타일 - 함수", [
        "함수는 짧고 한가지 일만 해야 한다.","",
        "함수에서 지역 변수를 10개 이하로 써야 한다.","",
        "함수 정의들 간에는 한줄을 띄워서 쓴다.","",
        "함수 선언에서 타입과 이름을 모두 쓴다.","",
    ]),
    ("C 코딩 스타일 - 함수 관련 goto", [
        "goto는 많은 프로그래머들이 더이상 쓰지 않는다.","",
        "여러 군데에서 함수가 종료되는데 공통된 정리 코드가 있을 때 유용하다.","",
        "int fun(int a)\n{\n        int result = 0;\n        char *buffer;\n\n        buffer = malloc(SIZE);\n        if (!buffer)\n                return -ENOMEM;\n\n        if (condition1) {\n                while (loop1) {\n                        ...\n                }\n                result = 1;\n                goto out_free_buffer;\n        }\n        ...\nout_free_buffer:\n        free(buffer);\n        return result;\n}",
    ]),
    ("C 코딩 스타일 - 함수 관련 goto 에러", [
        "one err bugs로 알려진 흔한 버그","",
        "err:\n        free(foo->bar);\n        free(foo);\n        return ret;",
        "foo가 NULL인 상태에서 err로 이동했을 때 버그 발생","",
    ]),
    ("C 코딩 스타일 - 함수 관련 goto 에러 해결책", [
        "err_free_bar:와 err_free_foo:로 라벨을 분리한다.","",
        "err_free_bar:\n        free(foo->bar);\nerr_free_foo:\n        free(foo);\n        return ret;",
        "err_free_bar로는 foo가 NULL이 아닐 때만 이동","",
    ]),
    ("C 코딩 스타일 - 주석", [
        "주석은 함수 앞에 달아서 무엇을 하고, 왜 하는지를 설명한다.","",
        "아주 복잡한 함수의 경우 본문 안에도 구획을 나눠서 주석을 달 수 있다.","",
        "/*\n * This is the preferred style for multi-line\n * comments in the Linux kernel source code.\n */",""
        "데이터에 대해 짧은 주석을 다는 것도 중요하다.",""
    ]),
    ("C 복습 - 이름의 범위", [
        "int x = 0;",
        "void function(void)",
        "{",
        "        int a = x;          /* 전역 x(=0) */",
        "        int x = 1;          /* 함수 지역 x */",
        "        for (int x = 2; x < 10; x++) {",
        "                int b = x;  /* for-블록의 x */",
        "        }",
        "        a = x;              /* 여기서의 x는 함수 지역 x(=1) */",
        "}","",
        "함수가 끝나기 전 a에 저장된 값은?",""
    ]),
    ("C 복습 - 배열.c", [
        "#include <stdio.h>","",
        "void array_test()\n{",
        "        int arr[3];",
        "        arr[0] = 1;",
        "        arr[1] = 2;",
        "        arr[2] = 3;",
        "        for (int i = 0; i < 3; i++)",
        '                printf("%d ", arr[i]);',"",
        '        printf("\\n");',
        "}","",
        "array_test()를 실행했을 때 화면에 출력되는 문자열은?"
    ]),
    ("C 복습 - 배열_포인터.c", [
        "#include <stdio.h>","",
        "void array_ptr_test()\n{",
        "        int arr[3];",
        "        int *arr_ptr = arr;",
        "        arr_ptr[0] = 1;",
        "        arr_ptr[1] = 2;",
        "        arr_ptr[2] = 3;",
        "        for (int i = 0; i < 3; i++)",
        '                printf("%d ", arr_ptr[i]);',"",
        '        printf("\\n");',
        "}","",
        "array_ptr_test()를 실행했을 때 화면에 출력되는 문자열은?"
    ]),
    ("C 복습 - 배열_포인터_파라미터.c", [
        "#include <stdio.h>","",
        "void array_set_elements(int *arr_ptr)\n{",
        "        arr_ptr[0] = 1;",
        "        arr_ptr[1] = 2;",
        "        arr_ptr[2] = 3;",
        "}","",
        "void array_fun_test()\n{",
        "        int arr[3];",
        "        array_set_elements(arr);",
        "        for (int i = 0; i < 3; i++)",
        '                printf("%d ", arr[i]);',"",
        '        printf("\\n");',
        "}","",
        "array_fun_test()를 실행했을 때 화면에 출력되는 문자열은?"
    ]),
    ("C 복습 - 구조체.c", [
        "#include <stdio.h>","",
        "struct Point {",
        "        int x;",
        "        int y;",
        "};","",
        "void struct_test()\n{",
        "        struct Point p;",
        "        p.x = 1;",
        "        p.y = 2;",
        '        printf("(%d, %d)\\n", p.x, p.y);',
        "}","",
        "struct_test()를 실행했을 때 화면에 출력되는 문자열은?"
    ]),
    ("C 복습 - 구조체_포인터.c", [
        "#include <stdio.h>","",
        "struct Point {",
        "        int x;",
        "        int y;",
        "};","",
        "void struct_ptr_test()\n{",
        "        struct Point p;",
        "        struct Point *p_ptr = &p;",
        "        p_ptr->x = 1;",
        "        p_ptr->y = 2;",
        '        printf("(%d, %d)\\n", p_ptr->x, p_ptr->y);',
        "}","",
        "struct_ptr_test()를 실행했을 때 화면에 출력되는 문자열은?"
    ]),
    ("C 복습 - 구조체_포인터_파라미터.c", [
        "#include <stdio.h>","",
        "struct Point {",
        "        int x;",
        "        int y;",
        "};","",
        "void struct_set_members(struct Point *p_ptr)\n{",
        "        p_ptr->x = 1;",
        "        p_ptr->y = 2;",
        "}","",
        "void struct_ptr_parameter_test()\n{",
        "        struct Point p;",
        "        struct_set_members(&p);",
        '        printf("(%d, %d)\\n", p.x, p.y);',
        "}","",
        "struct_ptr_parameter_test()를 실행했을 때 화면에 출력되는 문자열은?"
    ]),
    ("C 복습 - 메모리_할당_해제.c", [
        "#include <stdio.h>","",
        "#include <stdlib.h>","",
        "struct Point {",
        "        int x;",
        "        int y;",
        "};","",
        "struct Point *create_point()\n{",
        "        return malloc(sizeof(struct Point));",
        "}","",
        "void free_point(struct Point *p_ptr)\n{",
        "        free(p_ptr);",
        "}","",
        "void struct_set_members(struct Point *p_ptr)\n{",
        "        p_ptr->x = 1;",
        "        p_ptr->y = 2;",
        "}","",
        "void struct_ptr_malloc_test()\n{",
        "        struct Point* p_ptr = create_point();",
        "        if (!p_ptr) \n                return;",
        "        struct_set_members(p_ptr);",
        '        printf("(%d, %d)\\n", p_ptr->x, p_ptr->y);',
        "        free_point(p_ptr);"
        "}","",
        "struct_ptr_malloc_test()를 실행했을 때 화면에 출력되는 문자열은?"
    ]),

    ("C 키워드 _Bool", [
        "0 또는 1을 저장할 수 있는 타입","",
        "거짓 또는 참을 저장하기 위해 쓰임","",
    ]),

]
# 슬라이드 생성
for title, bullets in slides:
    add_slide(title, bullets)

# 저장
output_path = "ppt/자료구조_5주차_10회차_PPT.pptx"
prs.save(output_path)
output_path
