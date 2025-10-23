from pptx import __version__ as pptx_version
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

print(pptx_version)

# 새 프레젠테이션 생성
prs = Presentation("template.pptx")
prs.slides.add_slide(prs.slide_masters[0].slide_layouts[0])  # 제목 슬라이드 추가
# 제목 슬라이드 설정
title_slide = prs.slides[0]
title_slide.shapes.title.text = "자료구조 (Data Structure)\n6주차: 힙"


layout = prs.slide_masters[0].slide_layouts[1]  # 제목과 내용 레이아웃 선택

# 슬라이드 추가 함수
def add_slide(title, bullets):
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    content = slide.placeholders[1]
    content.text = bullets[0] if bullets else ""
    for bullet in bullets[1:]:
        p = content.text_frame.add_paragraph()
        p.text = bullet

slides = [
    ("이진 트리를 포인터로 구현하기 정리", [
        "이진 트리 노드 타입",
        "\t - left: \t왼쪽 자식 주소: ","",
        "\t - right: \t오른쪽 자식 주소: ","",
        "\t - data: \t  저장된 데이터: ","",
    ]),
    ("이번 시간 목차", [
        "이진 트리를 배열로 구현하기","",
        "특별한 속성을 가진 이진 트리","",
        "특별한 속성을 가진 이진 트리 예제","",
    ]),
    ("이진 트리를 배열로 구현하기 - 배열", [
        "배열의 특징은?","",
    ]),
    ("이진 트리를 배열로 구현하기 - 저장 순서", [
        "배열의 첫번째에 저장하기에 적합한 데이터는?","",
        "배열의 두번째에 저장하기에 적합한 데이터는?","",
    ]),
    ("이진 트리를 배열로 구현하기 - 순서 계산", [
        "자식 데이터의 순서를 계산하는 방법은?","",
        "부모 데이터의 순서를 계산하는 방법은?","",
    ]),
    ("이진 트리를 배열로 구현하기 - 추가, 삭제", [
        "데이터를 추가할 적절한 위치는?","",
        "데이터를 삭제할 때의 적절한 동작은?","",
    ]),
    ("이진 트리를 배열로 구현하기 - 구조체", [
        "배열로 구현한 이진 트리 타입","",
        "\t - data: \t배열 주소","",
        "\t - capacity: \t배열 크기","",
        "\t - size: \t현재 저장된 데이터 개수","",
    ]),
    ("특별한 속성을 가진 이진 트리", [
        "이진 트리에 어떻게 특별한 속성을 부여할 수 있을까?","",
    ]),
    ("특별한 속성을 가진 이진 트리", [
        "힙: 부모 노드와 자식 노드 간에 특정 관계가 항상 성립한다","",
        "\t - 최대 힙: 부모 노드가 제일 크다","",
        "\t - 최소 힙: 부모 노드가 제일 작다","",
        "\t - 이진 탐색 트리: 왼쪽 자식 < 부모 < 오른쪽 자식","",
    ]),
    ("특별한 속성을 가진 이진 트리 - 최대 힙", [
        "최대 힙에서 데이터를 추가할 적절한 위치는?","",
    ]),
    ("특별한 속성을 가진 이진 트리 - 최대 힙", [
        "함수 MAX_HEAP_INSERT(heap, value):",
        "실행조건: heap->size < heap->capacity",
        "++heap->size",
        "heap->data[heap->size] = value",
        "HEAPIFY_UP(heap, heap->size)",
    ]),
    ("특별한 속성을 가진 이진 트리 - 최대 힙", [
        "최대 힙의 끝에 데이터가 추가된 경우,",
        " 다시 힙 속성을 회복하려면?","",
    ]),
    ("특별한 속성을 가진 이진 트리 - 최대 힙", [
        "함수 HEAPIFY_UP(heap, child):",
        "만약: child == 0이면 반환",
        "parent = child / 2",
        "만약: heap->data[parent] < heap->data[child]이면",
        "\ttmp = heap->data[parent]",
        "\theap->data[parent] = heap->data[child]",
        "\theap->data[child] = tmp",
        "\tHEAPIFY_UP(heap, parent)",
    ]),
    ("특별한 속성을 가진 이진 트리 - 최대 힙", [
        "다음 최대 힙에 5, 3, 8, 1, 2를 순서대로 추가하면?",
        "",
        "        8",
        "       / \\",
        "      3   5",
        "     / \\",
        "    1   2",
    ]),
    ("특별한 속성을 가진 이진 트리 - 최대 힙", [
        "최대 힙에서 데이터를 삭제한 경우, 그 자리를 메꿀 노드는?","",
    ]),
    ("특별한 속성을 가진 이진 트리 - 최대 힙", [
        "함수 MAX_HEAP_DELETE(heap, index):",
        "실행조건: heap->size > 0",
        "value = heap->data[index]",
        "heap->data[index] = heap->data[heap->size]",
        "heap->size--",
        "HEAPIFY_DOWN(heap, index)",
        "반환: value",
    ]),
    ("특별한 속성을 가진 이진 트리 - 최대 힙", [
        "최대 힙에서 삭제된 데이터 자리를 제일 끝 노드가 메꿨다면, 힙 속성을 회복할 방법은?","",
    ]),
    ("특별한 속성을 가진 이진 트리 - 최대 힙", [
        "함수 HEAPIFY_DOWN(heap, parent):",
        "left = parent * 2",
        "right = parent * 2 + 1",
        "largest = parent",
        "만약: left <= heap->size",
        "\t&& heap->data[left] > heap->data[largest]이면",
        "\tlargest = left",
        "만약: right <= heap->size",
        "\t&&heap->data[right] > heap->data[largest]이면",
        "\tlargest = right",
        "만약: largest != parent이면",
        "\ttmp = heap->data[parent]",
        "\theap->data[parent] = heap->data[largest]",
        "\theap->data[largest] = tmp",
        "\tHEAPIFY_DOWN(heap, largest)",
    ]),

    ("특별한 속성을 가진 이진 트리 - 최대 힙", [
        "다음 최대 힙에서 루트 노드를 삭제하면?",
        "",
        "        8",
        "       / \\",
        "      3   5",
        "     / \\",
        "    1   2",
        "",
        "        5",
        "       / \\",
        "      3   2",
        "     /",
        "    1",
    ]),

    ("예제 - 우선순위 큐", [
        "우선순위가 높은 데이터가 먼저 삭제되는 큐","",
        "\t - 최대 힙을 이용하여 구현한다","",
        "\t - 추가: 최대 힙에 데이터 추가","",
        "\t - 삭제: 최대 힙에서 루트 노드를 제거","",
    ]),


    ("예제 - 소셜미디어 데이터", [
        "소셜미디어 홈페이지에서 관리해야 할 자료는?","",
    ]),

    ("예제 - 소셜미디어 데이터", [
        "각 유저마다 저장할 정보는?","",
    ]),

    ("예제 - 소셜미디어 데이터", [
        "유저 데이터 구조체","",
        "\t - 친구 유저 리스트","",
        "\t - 게시물 리스트 (최신순으로 정렬됨)","",
    ]),

    ("예제 - 소셜미디어 데이터", [
        "유저의 최신 게시물과 친구들의 최신 게시물들을",
        "통틀어 가장 최신인 10개 게시물 리스트를 만들려면?","",
    ]),

    ("예제 - 소셜미디어 데이터", [
        "함수 GET_TIMELINE(user):",
        "result = 새로운 게시물 리스트",
        "Q = 새로운 게시물 노드 우선순위 큐",

        "반복: user의 모든 친구 friend들에 대해",
        "\tfriend = GET_USER(friend_id)",
        "\t만약: friend의 게시물이 있으면",
        "\t\tENQUEUE(Q, friend의 게시물들의 first)",

        "반복: SIZE(Q) > 0 그리고 SIZE(result) < 10",
        "\tnode = DEQUEUE(Q)",
        "\tpush(result, node->data)",
        "\t만약: node->next != NULL이면",
        "\t\tENQUEUE(Q, node->next)",

        "반환: result",
    ]),

    ("예제 - 소셜미디어 데이터", [
        "참고 자료",
        "https://leetcode.com/problems/design-twitter/",
    ]),

    ("이번 시간 정리", [
        "이진 트리를 배열로 구현하기","",
        "힙의 정의와 데이터 추가, 삭제","",
        "힙으로 우선순위 큐 만들기","",
        "소셜 미디어 타임라인 예제","",
    ]),

    ("2021 중간고사 기출문제", [
        "아주대학교 재학생들의 소프트웨어학과 개설 과목",
        "이수 여부를 저장할 때 자료 구조를 추천하시오",
    ]),

    ("2021 중간고사 기출문제", [
        "매일 행적을 기록할 때 자료 구조를 추천하시오",
    ]),

    ("2021 중간고사 기출문제", [
        "식당에서 손님들의 주문 내역을 기록할 때 쓸",
        "자료 구조를 추천하시오",
    ]),    

    ("2021 중간고사 기출문제", [
        "여행 계획을 세울 때, 여러 상황 및 대응 계획을",
        "저장할 자료구조를 추천하시오 (다음 상황 중복 없음)",
    ]),    

    ("2021 중간고사 기출문제", [
        "쇼핑할 목록이 자주 바뀐다고 가정했을 때,",
        "쇼핑할 목록을 저장할 자료구조를 추천하시오",
    ]),    

    ("2021 중간고사 기출문제", [
        "Stack에 데이터를 추가, 삭제할 때 드는 시간은?",
    ]),    

    ("2021 중간고사 기출문제", [
        "Queue에 데이터를 추가, 삭제할 때 드는 시간은?",
    ]),    

    ("2021 중간고사 기출문제", [
        "Singly linked list의 시작 node 주소만 알고 있을 때,",
        "특정 값을 가진 node를 삭제할 때 드는 시간은?",
    ]),    

    ("2021 중간고사 기출문제", [
        "Doubly linked list의 head node 주소만 알 때,",
        "마지막 노드를 삭제하는데 드는 시간은?",
    ]),    

    ("2021 중간고사 기출문제", [
        "Binary Tree에서 i층에 있는 노드들의 data를",
        "출력하는데 드는 시간은?",
    ]),    

    ("2021 중간고사 기출문제", [
        "chain 형태의 linked list의 데이터들을 리스트 내에",
        "저장되어 있는 순서와 반대로 출력하는 재귀 함수는?",
    ]),    

    ("2021 중간고사 기출문제", [
        "1 + 2 형태로 표현된 수식을 1 2 + 형태로 변환하는",
        "코드를 쓰시오",
    ]),

    ("2021 중간고사 기출문제", [
        "chain 형태의 linked list에 노드를 추가하는 코드를",
        "쓰시오",
    ]),

]

# 슬라이드 생성
for title, bullets in slides:
    add_slide(title, bullets)


def add_table_slides(title, bullets, table_data):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title

        content = slide.placeholders[1]
        content.text = bullets[0] if bullets else ""
        for bullet in bullets[1:]:
                p = content.text_frame.add_paragraph()
                p.text = bullet

        cols, rows = len(table_data), len(table_data[0]) if table_data else 0
        left = Inches(0.5)
        top = Inches(5)
        width = Inches(9)
        height = Inches(0.8 + 0.3 * rows)
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        for c in range(cols):
                for r in range(rows):
                        table.cell(r, c).text = table_data[c][r]


table_slides = [
        ("트리 방문 - Preorder", [
                "1    만약: N == NULL이면 반환",
                "2    pre_task(N)",
                "3    PREORDER(LEFT(N))",
                "4    PREORDER(RIGHT(N))",
        ], [
                ["호출스택", "pre_task", "줄번호"],
                ["A", "A", "3"],
                ["AB", "B", "3"],
                ["ABn", "", "1"],
                ["AB", "", "4"],
                ["AB", "", "-"],
                ["A", "", "4"],
                ["AC", "C", "3"],
                ["ACn", "", "1"],
                ["AC", "", "4"],
                ["ACn", "", "1"],
                ["AC", "", "4"],
                ["A", "", "4"],
                ["An", "", "1"],
        ]),
]

for title, bullets, table_data in table_slides:
        add_table_slides(title, bullets, table_data)


# 저장
output_path = "ppt/자료구조_6주차_11회차_PPT.pptx"
prs.save(output_path)
output_path
