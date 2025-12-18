from pptx import __version__ as pptx_version
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

print(pptx_version)

# 새 프레젠테이션 생성
prs = Presentation("template.pptx")
prs.slides.add_slide(prs.slide_masters[0].slide_layouts[0])  # 제목 슬라이드 추가
# 제목 슬라이드 설정
title_slide = prs.slides[0]
title_slide.shapes.title.text = "자료구조 (Data Structure)\n7주차: 정렬"


layout = prs.slide_masters[0].slide_layouts[1]  # 제목과 내용 레이아웃 선택

# 슬라이드 추가 함수
def add_slide(title, bullets):
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    content = slide.placeholders[1]
    content.text = bullets[0] if bullets else ""
    for bullet in bullets[1:]:
        p = content.text_frame.add_paragraph()
        p.text = bullet

slides = [
    ("오늘의 목차", [
        "정렬 알고리즘의 필요성과 정의","",
        "정렬 알고리즘의 종류","",
        "\t heap sort와 selection sort","",
        "\t quick sort와 insertion sort","",
        "\t merge sort와 hybrid sort","",
        "\t counting sort와 radix sort","",
    ]),
    ("정렬 알고리즘의 필요성", [
        "데이터를 정렬하는 이유는?","",
        "\t - 유저의 필요","",
        "\t - 프로그램 성능 향상","",
    ]),
    ("정렬의 정의", [
        "데이터를 특정 기준에 따라 순서대로 배열하는 과정","",
        "오름차순과 내림차순 정렬","",
    ]),
    ("HEAP_SORT", [
        "힙 자료구조를 이용한 정렬 방법","",
    ]),
    ("HEAP_SORT 과정", [
        "1. 배열로부터 최대 힙 생성","",
        "2. 루트 노드를 제거하고 마지막 노드와 교체","",
        "3. 나머지 노드들에 대해 힙 속성 회복","",
        "4. 2~3 과정을 반복하여 정렬 완료","",
    ]),
    ("HEAP_SORT 시각화", [
        "예시: [5, 3, 2, 7, 1] 정렬 과정","",
        "\t - 최대 힙 생성: [7, 5, 2, 3, 1]",
        "\t - 1회전: [5, 3, 2, 1] + [7]",
        "\t - 2회전: [3, 1, 2] + [5, 7]",
        "\t - 3회전: [2, 1] + [3, 5, 7]",
        "\t - 4회전: [1] + [2, 3, 5, 7]",
        "\t - 5회전: [] + [1, 2, 3, 5, 7]",
    ]),
    ("HEAP_SORT", [
        "함수 HEAP_SORT(배열리스트 L):","",
        "힙 heap = 새 최대힙(L->data, L->size)","",
        "반복: i는 L의 전체 인덱스를 거꾸로","",
        "\t L->data[i] = HEAP_POP(heap)","",
        "삭제: heap","",
    ]),
    ("HEAP_SORT의 시간 복잡도", [
        "최대 힙 생성시 시간복잡도 (perfect binary tree): ","",
        "\t - 트리의 가장 끝 층(leaf)의 노드들: 비교 0회, 교환 0회",
        "\t - 끝에서 두번째 층(leaf 위)의 노드들: 비교 2회, 교환 1회",
        "\t - ... 루트 노드: 비교 2h 회, 교환 h 회","",
        "\t - 비교 총합: n / 2 * 0 + n / 2^1 * 1 + ... + n / 2^h * h","",
        "\t\t\t = n * (sigma h / 2^h)","",
        "\t\t\t ≒ 2n","",
        "\t - 교환 총합 ≒ n",
        "\t - sigma x^h = 1 / (1 - x)",
        "\t - sigma h * x^(h-1) = 1 / (1 - x)^2","",
        "\t - sigma h * x^h = x / (1 - x)^2","",
    ]),
    ("HEAP_SORT의 시간 복잡도 (계속)", [
        "HEAP_POP 수행시 시간복잡도 (perfect binary tree):","",
        "\t - 트리 높이가 h이면: 비교 2h회, 교환 h회","",
        "\t - ... 높이가 1이면: 비교 2회, 교환 1회","",
        "\t - 비교 총합: 2 * h * n / 2^1 + 2 * (h - 1) * n / 2^2 + ...","",
        "\t\t\t = sigma (i = 0 to h-1) ((h - i) * n / 2^i)","",
        "\t\t\t = hn sigma (i = 0 to h-1) (1 / 2^i) - n sigma (i = 0 to h-1) (i / 2^i)","",
        "\t\t\t = 2 * h * n - 2 * n",
        "\t\t\t ≒ 2n log n","",
    ]),

    ("SELECTION_SORT", [
        "HEAP_SORT와 같은 전략",
        "\t - 정렬 전 구간의 최대값을 맨 뒤로 이동",
        "최대값을 찾을 때 모든 값을 확인해서 느림","",
    ]),
    ("SELECTION_SORT 과정", [
        "1. 배열에서 최대값 찾기","",
        "2. 최대값을 맨 뒤의 원소와 교체","",
        "3. 두번째 최대값을 뒤에서 두번째 원소와 교체","",
        "4. ... 반복"
    ]),
    ("SELECTION_SORT 시각화", [
        "예시: [5, 3, 2, 7, 1] 정렬 과정","",
        "1회전: [5, 3, 2, 1, 7]","",
        "2회전: [1, 3, 2, 5, 7]","",
        "3회전: [1, 2, 3, 5, 7]","",
        "4회전: [1, 2, 3, 5, 7]","",
        "5회전: [1, 2, 3, 5, 7]","",
    ]),
    ("SELECTION_SORT", [
        "함수 SELECTION_SORT(리스트 L):","",
        "반복: pos는 마지막 인덱스부터 시작 인덱스까지","",
        "\t max_i = pos","",
        "\t 반복: i는 pos부터 시작 인덱스까지","",
        "\t\t 만약: L->data[i] > L->data[max_i]이면","",
        "\t\t\t max_i = i","",
        "\t SWAP(L, max_i, pos)","",
    ]),
    ("SELECTION_SORT", [
        "함수 SWAP(리스트 L, 인덱스 i, 인덱스 j):","",
        "tmp = L->data[i]","",
        "L->data[i] = L->data[j]","",
        "L->data[j] = tmp","",
    ]),
    ("SELECTION_SORT의 시간 복잡도", [
        "비교 횟수:",
        "최대값을 찾을 때: n - 1회 비교",
        "두번째 최대값을 찾을 때: n - 2회 비교",
        "...",
        "총합: (n - 1) + (n - 2) + ... + 1 = n(n - 1) / 2 : O(n^2)","",
        "교환 횟수: n - 1회 : O(n)","",
    ]),
    ("HEAP_SORT와 SELECTION_SORT 비교", [
        "시간 복잡도",
        "\t - HEAP_SORT: O(n log n)",
        "\t - SELECTION_SORT: O(n^2)","",
        "공간 복잡도",
        "\t - 둘 다 O(1)","",
        "힙 정렬은 빅데이터에서 더 나은 성능을 제공","",
    ]),





    ("QUICK_SORT", [
        "가장 빠른 정렬 알고리즘","",
        "기준값을 정해서 작은 값들과 큰 값들을 분리, 재배치","",
    ]),

    ("QUICK_SORT 과정", [
        "1. 기준값 선택 (제일 끝 값)","",
        "2. 기준값을 기준으로 작은 값들과 큰 값들로 분할","",
        "3. 작은 값들과 큰 값들에 대해 1~2 과정을 반복","",
    ]),

    ("QUICK_SORT 시각화", [
        "예시: [5, 3, 2, 7, 1] 정렬 과정","",
        "기준값 5 선택: [3, 2, 1] ,\t [5] ,\t\t [7]",
        "기준값 3 선택: [2, 1] , \t[3]",
        "기준값 2 선택: [1] , [2]","",
        "최종 정렬 결과: [1, 2, 3, 5, 7]","",
    ]),

    ("QUICK_SORT", [
        "함수 QUICK_SORT(리스트 L, 인덱스 lo, 인덱스 hi):","",
        "만약: lo < hi이면","",
        "\t pivot = PARTITION(L, lo, hi)","",
        "\t QUICK_SORT(L, lo, pivot - 1)","",
        "\t QUICK_SORT(L, pivot + 1, hi)","",
    ]),

    ("QUICK_SORT", [
        "함수 PARTITION(리스트 L, 인덱스 lo, 인덱스 hi):",
        "pivot_data = L->data[hi]",
        "left_rear = lo \t // pivot이하 그룹 끝 + 1",
        "반복: i는 lo부터 hi - 1까지",
        "\t 만약: L->data[i] <= pivot_data이면",
        "\t\tSWAP(L, left_rear, i)",
        "\t\tleft_rear += 1",
        "SWAP(L, left_rear, hi)",
        "반환: left_rear",
    ]),

    ("QUICK_SORT의 시간 복잡도", [
        "n개를 QUICK_SORT할 때 비교 횟수 (왼쪽이 k개일 때):","",
        "C(n) = n - 1 + C(k) + C(n - k - 1)","",
        "\t - k는 피벗보다 작은 원소의 개수","",
        "\t - 평균적으로, C(n) = n - 1 + 2 * sigma k=0~n-1 C(k) / n","",
        "\t - 다시 쓰면, n * C(n) = n * (n - 1) + 2 * sigma k=0~n-1 C(k)","",
        "\t - 또한, (n - 1) * C(n - 1) = (n - 1) * (n - 2) + 2 * sigma k=0~n-2 C(k)","",
        "\t - 위 두 식을 빼면, n * C(n) - (n - 1) * C(n - 1) = (n - 1)(n - 2) + 2 * C(n - 1)","",
        "\t - 따라서, n * C(n) = (n + 1) * C(n - 1) + 2(n - 1)","",
        "\t - 다시 쓰면, C(n) / (n + 1) = C(n - 1) / n + 2(n - 1) / n(n + 1)","",
        "\t\t\t  = C(n) / (n + 1) + 2 / n - 4 / (n(n + 1))",
        "\t\t\t  ~~ 2 sigma 1 / n",
        "\t - 적분으로 근사하면, 2 sigma 1 / n ~~ 2 ln n","",
        "\t - 따라서, C(n) ~~ 2n ln n","",
        "평균 교환 횟수는 C(n)/2 ~~ n ln n","",
    ]),

    ("QUICK_SORT의 시간 복잡도", [
        "최악의 경우 비교 횟수 (왼쪽이 0개로 나뉠 때):","",
        "\t - 이미 정렬된 배열에서 피벗이 가장 크거나 작은 경우","",
        "\t - 비교 횟수 C(n) = n - 1 + C(n - 1)","",
        "\t - 따라서, C(n) = (n - 1) + (n - 2) + ... + 1 = n(n - 1) / 2 : O(n^2)","",
    ]),

    ("INSERTION_SORT", [
        "정렬된 부분의 적절한 위치로 데이터를 이동해간다","",
    ]),

    ("INSERTION_SORT 과정", [
        "왼쪽부터 정렬해나가는 경우:",
        "1. 정렬되지 않은 부분의 왼쪽 원소에 대해서,\n\t 정렬된 부분에서 적절한 위치 찾기","",
        "2. 찾은 위치에 추가하여 정렬된 부분 확장","",
        "3. ... 반복","",
    ]),
    ("INSERTION_SORT 시각화", [
        "예시: [5, 3, 2, 7, 1] 정렬 과정","",
        "\t 1회전: [3, 5, 2, 1, 7]","",
        "\t 2회전: [2, 3, 5, 1, 7]","",
        "\t 3회전: [1, 2, 3, 5, 7]","",
    ]),
    ("INSERTION_SORT", [
        "함수 INSERTION_SORT(리스트 L, 인덱스 lo, 인덱스 hi):",
        "반복: rear는 lo부터 hi까지",
        "\t value = L->data[rear] \t // 이동할 데이터",
        "\t i = rear \t\t\t // 데이터를 이동할 위치",
        "\t 반복: i > lo이고, L->data[i - 1] > value인 동안",
        "\t\t L->data[i] = L->data[i - 1]",
        "\t\t i -= 1",
        "\t L->data[i] = value",
    ]),
    ("INSERTION_SORT의 시간 복잡도", [
        "반복 조건 확인 횟수 (최악):",
        "C_worst(n) = 1 + 2 + ... + n = n(n + 1) / 2 : O(n^2)",
        "반복 조건 확인 횟수 (평균):",
        "C_avg(n) = C_worst(n) / 2 = n(n + 1) / 4 : O(n^2)",
        "반복 조건 확인 횟수 (최선):",
        "C_best(n) = 1 + 1 + ... + 1 = n : O(n)",        
    ]),


    ("퀵 정렬과 힙 정렬 속도 비교", [
        "퀵 정렬이 단순해서 평균 실행 시간이 더 빠르다","",
        "힙 정렬은 최악의 경우에도 O(n log n)을 보장한다","",
    ]),

    ("하이브리드 정렬 (INTRO_SORT)", [
        "퀵 정렬의 빠른 평균 성능과 힙 정렬의 최악 성능 결합","",
        "n이 작아지면 INSERTION_SORT로 전환","",
        "C++의 표준 정렬 함수","",
    ]),
    ("INTRO_SORT 과정", [
        "1. 퀵 정렬로 분할 정복 수행","",
        "2. 재귀 깊이가 일정 수준 이상일 경우 힙 정렬로 전환","",
        "3. 데이터 수가 적으면 INSERTION SORT","",
    ]),
    ("INTRO_SORT 시각화", [
        "예시: [5, 3, 2, 7, 1] 정렬 과정","",
        "삽입 정렬이 실행됨","",
        "최종 정렬 결과: [1, 2, 3, 5, 7]","",
    ]),
    ("INTRO_SORT", [
        "함수 INTRO_SORT(리스트 list, 인덱스 lo, 인덱스 hi, 정수 limit):","",
        "만약: hi - lo < 15이면","",
        "\t INSERTION_SORT(list, lo, hi)","",
        "아니면 만약: limit == 0이면","",
        "\t HEAP_SORT(list, lo, hi)","",
        "아니면","",
        "\t pivot = PARTITION(list, lo, hi)","",
        "\t INTRO_SORT(list, lo, pivot - 1, limit - 1)","",
        "\t INTRO_SORT(list, pivot + 1, hi, limit - 1)","",
    ]),


    ("MERGE SORT", [
        "반으로 나눠서 각각 정렬한 후, 정렬하며 합치는 방법","",
    ]),
    ("MERGE SORT 과정", [
        "1. 배열을 반으로 분할","",
        "2. 각 부분에 대해 재귀적으로 정렬","",
        "3. 정렬된 두 부분을 합침","",
    ]),
    ("MERGE SORT 시각화", [
        "예시: [5, 3, 2, 7, 1] 정렬 과정","",
        "분할: [5,3] , [2,7,1]","",
        "병합: [3,5] , [1,2,7]","",
        "최종 정렬 결과: [1,2,3,5,7]","",
    ]),
    ("MERGE_SORT", [
        "함수 MERGE_SORT(리스트 L, 인덱스 lo, 인덱스 hi):","",
        "만약: lo < hi이면",
        "\t mid = (lo + hi) / 2",
        "\t MERGE_SORT(L, lo, mid)",
        "\t MERGE_SORT(L, mid + 1, hi)",
        "\t MERGE(L, lo, mid, hi)",
    ]),
    ("MERGE", [
        "함수 MERGE(리스트 L, 인덱스 lo, 인덱스 mid, 인덱스 hi):","",
        "left_size = mid - lo + 1",
        "right_size = hi - mid","",
        "left_q = 새로운 큐(left_size)","",
        "right_q = 새로운 큐(right_size)","",
        "반복: i는 lo부터 mid까지","",
        "\t ENQUEUE(left_q, L->data[i])","",
        "반복: i는 mid + 1부터 hi까지","",
        "\t ENQUEUE(right_q, L->data[i])","",
        "merged = 새로운 리스트(L->data, 0)","",
        "반복: !EMPTY(left_q) 그리고 !EMPTY(right_q)일 동안","",
        "\t 만약: FRONT(left_q) <= FRONT(right_q)이면","",
        "\t\t ALIST_PUSH(merged, DEQUEUE(left_q))","",
        "\t 아니면","",
        "\t\t ALIST_PUSH(merged, DEQUEUE(right_q))","",
        "반복: !EMPTY(left_q)일 동안","",
        "\t ALIST_PUSH(merged, DEQUEUE(left_q))","",
        "반복: !EMPTY(right_q)일 동안","",
        "\t ALIST_PUSH(merged, DEQUEUE(right_q))","",
        "삭제: left_q, right_q, merged","",
    ]),
    ("MERGE_SORT의 시간 복잡도", [
        "MERGE시 시간복잡도:",
        "총 길이가 k인 두 배열을 합치는데 비교 k회, 대입 2k회",
        "MERGE 트리의 같은 층에 있는 총 작업량: 비교 n회, 대입 2n회",
        "MERGE 트리의 높이인 log n번 반복됨",
        "총 비교 n log n 회, 대입 2n log n회 : O(n log n)","",
    ]),

    ("하이브리드 정렬(Powersort)", [
        "앞에서부터 정렬된 구간을 찾고, 그 구간들을 MERGE하는 방법","",
        "정렬된 구간이 짧으면 INSERTION SORT로 늘림","",
        "정렬된 구간들은 스택에 보관","",
        "MERGE 트리에서 예상 깊이가 작거나 같은 구간들은 MERGE 수행","",
        "마지막에 스택의 모든 구간 MERGE","",
        "파이썬의 정렬 함수에서 쓰임","",
    ]),


    ("COUNTING_SORT", [
        "데이터의 개수를 세어 정렬하는 방법","",
        "시간 복잡도: O(n + k)","",
    ]),
    ("COUNTING_SORT 과정", [
        "1. 데이터 개수를 세는 배열 생성 및 초기화","",
        "2. 각 원소의 개수 세기","",
        "3. 출력 배열에 정렬된 원소 배치","",
    ]),
    ("COUNTING_SORT 시각화", [
        "예시: [5, 5, 2, 2, 1] 정렬 과정","",
        "최대값: 5","",
        "카운트 배열: [0, 1, 2, 0, 0, 2]","",
        "최종 정렬 결과: [1, 2, 2, 5, 5]","",
    ]),
    ("COUNTING_SORT", [
        "함수 COUNTING_SORT(리스트 L, 정수 max_value):","",
        "count = 새로운 배열(max_value + 1) // 0으로 초기화","",
        "반복: value는 L의 모든 데이터","",
        "\t count[value] += 1","",
        "sorted = 새로운 리스트(L->data, 0)","",
        "반복: value는 0부터 max_value까지 증가","",
        "\t 반복: count[value]번","",
        "\t\t ALIST_PUSH(value)","",
        "삭제: count, sorted","",
    ]),
    ("COUNTING_SORT의 시간 복잡도", [
        "개수 배열 생성 및 초기화: O(k)","",
        "데이터 개수 세기: O(n)","",
        "총 O(n + k)","",
    ]),


    ("RADIX_SORT", [
        "자리수별로 정렬하는 방법","",
    ]),
    ("RADIX_SORT 과정", [
        "1. 가장 낮은 자리수부터 시작","",
        "2. 자리수별로 안정 정렬 수행","",
        "3. 모든 자리수에 대해 반복","",
    ]),
    ("RADIX_SORT 시각화", [
        "예시: [170, 45, 75, 90, 802, 24, 2, 66] 정렬 과정","",
        "1의 자리수 정렬: [170, 90, 802, 2, 24, 45, 75, 66]","",
        "10의 자리수 정렬: [802, 2, 24, 45, 66, 75, 90, 170]","",
        "100의 자리수 정렬: [2, 24, 45, 66, 75, 90, 170, 802]","",
    ]),
    ("RADIX_SORT", [
        "함수 RADIX_SORT(리스트 list, 정수 max_value):","",
        "max_digits = 자리수 계산(max_value)","",
        "반복: digit는 0부터 max_digits - 1까지 증가","",
        "\t COUNTING_SORT_BY_DIGIT(list, digit)","",
    ]),
    ("COUNTING_SORT_BY_DIGIT", [
        "함수 COUNTING_SORT_BY_DIGIT(리스트 list, 정수 digit):","",
        "count_array = 새로운 리스트(10) // 0으로 초기화","",
        "반복: i는 list의 시작 인덱스부터 마지막 인덱스까지 증가","",
        "\t digit_value = (list->data[i] / 10^digit) % 10","",
        "\t count_array->data[digit_value]++","",
        "i = 0","",
        "반복: value는 0부터 9까지 증가","",
        "\t 반복: count는 1부터 count_array->data[value]까지 증가","",
        "\t\t list->data[i] = value","",
        "\t\t i++","",
        "삭제: count_array 해제","",
    ]),

   ("요약", [
        "heap sort: O(n log n)","",
        "quick sort: 평균 O(n log n), 최악 O(n^2)","",
        "insertion sort: 최선 O(n), 평균 및 최악 O(n^2)","",
        "merge sort: O(n log n)","",
        "counting sort: O(n + k)","",
   ]),
]

# 슬라이드 생성
for title, bullets in slides:
    add_slide(title, bullets)


output_path = "ppt/자료구조_7주차_13회차_PPT.pptx"
prs.save(output_path)
output_path
