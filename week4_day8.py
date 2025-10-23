from pptx import __version__ as pptx_version
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

print(pptx_version)

# 새 프레젠테이션 생성
prs = Presentation("template.pptx")
prs.slides.add_slide(prs.slide_masters[0].slide_layouts[0])  # 제목 슬라이드 추가
# 제목 슬라이드 설정
title_slide = prs.slides[0]
title_slide.shapes.title.text = "자료구조 (Data Structure)\n4주차: 큐"


layout = prs.slide_masters[0].slide_layouts[1]  # 제목과 내용 레이아웃 선택

# 슬라이드 추가 함수
def add_slide(title, bullets):
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    content = slide.placeholders[1]
    content.text = bullets[0] if bullets else ""
    for bullet in bullets[1:]:
        p = content.text_frame.add_paragraph()
        p.text = bullet

slides = [

    ("큐", [
        "First In First Out (FIFO) 자료구조","",
        "가장 먼저 들어온 데이터가 가장 먼저 나간다","",
        "데이터 추가/삭제는 반대쪽 끝에서 일어난다","",
    ]),

    ("오늘 목표", [
        "단방향 링크드 리스트로 구현한 큐","",
        "배열로 구현한 큐","",
        "미로 찾기","",
    ]),

    ("필수 헤더", [
        "#include <stdio.h>","",
        "#include <stdlib.h>","",
    ]),

    ("링크드 리스트-큐 구조체", [
        "더미 노드를 써서 가장 오래된 데이터 노드를 항상 front의 next에 두면, front가 NULL인 경우가 없다","",
        'front: 최근에 dequeue된 노드 (초기값 더미 노드)',
        "rear: 최근에 추가된 노드 (초기값 더미 노드)",
        "size: 데이터 개수 (초기값 0)","",
        "LQUEUE_ENQUEUE(q, value): 새 노드를 rear 뒤에 연결",
        "LQUEUE_DEQUEUE(q): 이전 front 삭제 / front의 데이터 반환",
    ]),

    ("링크드 리스트-큐 구조체 C 코드", [
        "struct lqueue {",
        "        struct node *front;",
        "        struct node *rear;",
        "        int size;",
        "};",
    ]),

    ("단방향 노드 개요", [
        "데이터와 다음 데이터 정보의 묶음","",
        "next: 다음 노드 정보",
        "value: 저장된 데이터","",
        "NODE_CREATE: 노드 생성하는 함수",
    ]),

    ("단방향 노드 C 코드", [
        "struct node {",
        "        struct node *next;",
        "        int value;",
        "};",
    ]),

    ("NODE_CREATE(value, next) C 코드", [
        "struct node *node_create(int value, struct node *next)",
        "{",
        "        struct node *node = malloc(sizeof(struct node));",
        "        if (!node) exit(1);",
        "        node->value = value;",
        "        node->next = next;",
        "        return node;",
        "}",
    ]),

    ("LQUEUE_CREATE() 개요", [
        "비어 있는 큐를 만들어서 반환한다","",
        "동작: q에 새로운 LQUEUE를 저장한다",
        "               dummy에 NODE_CREATE(0, NULL)을 저장한다",
        "               q의 front에 dummy를 저장한다",
        "               q의 rear에 dummy를 저장한다",
        "               q의 size에 0을 저장한다",
        "반환: q","",
    ]),

    ("LQUEUE_CREATE() C 코드", [
        "struct lqueue *lqueue_create()",
        "{",
        "        struct lqueue *q = malloc(sizeof(struct lqueue));",
        "        if (!q) exit(1);",
        "        q->front = q->rear = node_create(0, NULL); // 더미",
        "        q->size = 0;",
        "        return q;",
        "}",
    ]),

    ("LQUEUE_ENQUEUE(q, value)", [
        "데이터가 value인 노드를 만들어 rear 뒤에 추가한다","",
        "조건: q가 실제 큐를 가리킨다","",
        "동작: node에 NODE_CREATE(value, NULL)를 저장한다",
        "               q의 rear의 next에 node를 저장한다",
        "               q의 rear에 node를 저장한다",
        "               q의 size를 1 증가시킨다",
    ]),

    ("LQUEUE_ENQUEUE(q, value) C 코드", [
        "void lqueue_enqueue(struct lqueue *q, int value)",
        "{",
        "        if (!q) exit(1);",
        "        q->rear->next = node_create(value, NULL);",
        "        q->rear = q->rear->next;",
        "        ++q->size;",
        "}",
    ]),

    ("LQUEUE_DEQUEUE(q)", [
        "첫 데이터 노드를 더미로 돌리고 데이터를 반환한다. ","",
        "조건: q의 front의 next 존재","",
        "동작: old_head에 q의 front를 저장한다",
        "               q의 front에 old_head의 next를 저장한다",
        "               old_head를 삭제한다",
        "               q의 size를 1 감소시킨다",
        "반환: q의 front의 value",
    ]),

    ("LQUEUE_DEQUEUE(q) C 코드", [
        "int lqueue_dequeue(struct lqueue *q)",
        "{",
        "        if (!q) exit(1);",

        "        struct node *old_head = q->front;",
        "        if (!old_head->next) exit(1);",

        "        q->front = old_head->next;",

        "        q->size--;",

        "        free(old_head);",

        "        return q->front->value;",
        "}",
    ]),

    ("LQUEUE 요약", [
        "LQUEUE_CREATE	: front와 rear를 더미로","",
        "LQUEUE_ENQUEUE	: rear를 새로운 노드로","",
        "LQUEUE_DEQUEUE	: front를 한칸 다음으로","",
    ]),

    ("배열로 구현한 큐 구조체", [
        "배열 인덱스의 범위는 1..capacity","",
        "data: 데이터를 저장할 배열 (크기는 capacity + 1)",
        "capacity: 배열의 최대 저장량",
        "front: 마지막에 삭제한 데이터의 인덱스 (초기값 1)",
        "rear: 마지막에 추가한 데이터의 인덱스 (초기값 0)",
        "size: 저장된 데이터 개수 (초기값 0)",
    ]),

    ("배열로 구현한 큐 구조체", [
        "struct aqueue {",
        "        int *data;",
        "        int capacity;",
        "        int front;",
        "        int rear;",
        "        int size;",
        "};",
    ]),
    ("AQUEUE_CREATE(cap)", [
        "struct aqueue *aqueue_create(int cap)",
        "{",
        "        struct aqueue *q = malloc(sizeof(struct aqueue));",
        "        if (!q) exit(1);",
        "        q->data = malloc(sizeof(int)*(cap + 1));",
        "        if (!q->data) exit(1);",
        "        q->capacity = cap;",
        "        q->front = 0;",
        "        q->rear = 0;",
        "        q->size = 0;",
        "        return q;",
        "}",
    ]),

    ("SIMPLE_AQUEUE_ENQUEUE(q, value)", [
        "q의 끝에 데이터 value를 추가한다","",
        "조건: q의 rear < q의 capacity","",
        "동작: q의 rear를 1 증가시킨다",
        "               q의 data[q의 rear]에 value를 저장한다",
        "               q의 size를 1 증가시킨다",
    ]),

    ("SIMPLE_AQUEUE_ENQUEUE(q, v)", [
        "void simple_aqueue_enqueue(struct aqueue *q, int v)",
        "{",
        "        if (!q) exit(1);",
        "        if (q->rear == q->capacity) exit(1);",
        "        q->data[++q->rear] = v;",
        "        ++q->size;",
        "}",
    ]),

    ("SIMPLE_AQUEUE_DEQUEUE(q)", [
        "q의 첫번째 데이터를 삭제하고 데이터를 반환한다","",
        "조건: q의 front ≤ q의 rear","",
        "동작: tmp에 q의 data[q의 front]를 저장한다",
        "               q의 front를 1 증가시킨다",
        "               q의 size를 1 감소시킨다",
        "반환: tmp",
    ]),

    ("SIMPLE_AQUEUE_DEQUEUE(q)", [
        "int simple_aqueue_dequeue(struct aqueue *q)",
        "{",
        "        if (!q) exit(1);",
        "        if (q->front == q->rear) exit(1);",
        "        int tmp = q->data[++q->front];",
        "        --q->size;",
        "        return tmp;",
        "}",
    ]),

    ("SIMPLE_AQUEUE 요약", [
        "AQUEUE_CREATE			: front = rear = 0","",
        "SIMPLE_AQUEUE_ENQUEUE	: ++rear","",
        "SIMPLE_AQUEUE_DEQUEUE	: ++front","",
    ]),

    ("원형 배열로 구현한 큐", [
        "원형 배열: 인덱스가 capacity를 넘어가면",
        "                       다시 1이 된다","",
        "size로 비어있음과 꽉차있음을 체크한다"
    ]),

    ("CIRCULAR_AQUEUE_ENQUEUE(q, value)", [
        "q의 끝에 데이터 value를 추가한다","",
        "조건: q의 size < q의 capacity","",
        "동작: q의 rear를 1 증가시킨다",
        "               만약: q의 rear가 q의 capacity를 넘어갔으면",
        "                           q의 rear에 1을 저장한다",
        "               q의 data[q의 rear]에 value를 저장한다",
        "               q의 size를 1 증가시킨다",
    ]),

    ("CIRCULAR_AQUEUE_ENQUEUE(q, value)", [
        "void circular_aqueue_enqueue(struct aqueue *q, int v)",
        "{",
        "        if (!q) exit(1);","",
        "        if (q->size == q->capacity) exit(1);","",
        "        ++q->rear;",
        "        if (q->rear > q->capacity)",
        "                q->rear = 1;",
        "        q->data[q->rear] = v;",
        "        ++q->size;",
        "}",
    ]),

    ("CIRCULAR_AQUEUE_DEQUEUE(q)", [
        "q의 front 데이터를 삭제하고 반환한다","",
        "조건: q의 size > 0","",
        "동작: q의 front를 1 증가시킨다",
        "               만약: q의 front가 q의 capacity를 넘어갔으면",
        "                           q의 front에 1을 저장한다",
        "               tmp에 q의 data[q의 front]를 저장한다",
        "               q의 size를 1 감소시킨다",
        "반환: tmp",
    ]),

    ("CIRCULAR_AQUEUE_DEQUEUE(q)", [
        "int circular_aqueue_dequeue(struct aqueue *q)",
        "{",
        "        if (!q || q->size == 0) exit(1);","",
        "        ++q->front;",
        "        if (q->front > q->capacity)",
        "                q->front = 1;",
        "        int tmp = q->data[q->front];",
        "        q->size--;",
        "        return tmp;",
        "}",
    ]),

    ("미로 정보: 2차원 배열 (배열의 배열)", [
        "미로[i][j]: (행 i, 열 j) 위치의 상태","",
        "0 : 통로(아직 처리하지 않음)","",
        "1 : 벽","",
        "row*len + col : 최단경로에서 이전 좌표의 선형 index","",
        "0행과 0열은 모두 벽(1)로 두어, 내부 좌표에서 0/1과의 값 충돌을 방지함","",
    ]),

    ("좌표 구조체", [
        "row: 행 좌표","",
        "col: 열 좌표","",
        "distance: 출발점부터의 거리 (출발점은 0)","",
    ]),

    ("좌표 구조체 C 코드", [
        "struct point {",
        "        int row;",
        "        int col;",
        "        int distance;",
        "};",
    ]),

    ("MAZE_STATUS(maze, point)", [
        "반환: maze[point의 row][point의 col]","",
        " 0 : 미방문 통로",
        " 1 : 벽",
        " ≥ len+1 : 최단거리 계산됨 - 이전 좌표의 선형 index",
        "-1 : 출발점 표시",
    ]),

    ("MAZE_STATUS(maze, point) C 코드", [
        "int maze_status(int **maze, struct point point)",
        "{",
        "        return maze[point.row][point.col];",
        "}",
    ]),

    ("MAZE_MARK(maze, point, from)", [
        "maze[point의 row][point의 col]에 from 저장","",
        "최단경로에서 이전 좌표가 (i, j)일 때 from == i * len + j",
        "// 일차적으로 처리가 됐다는 정보가 저장됨",
    ]),

    ("MAZE_MARK(maze, point, from) C 코드", [
        "void maze_mark(int **maze, struct point point, int val)",
        "{",
        "        maze[point.row][point.col] = val;",
        "}",
    ]),

    ("좌표 선형 큐", [
        "points: 좌표 구조체의 배열","",
        "front: 마지막으로 꺼낸 좌표의 인덱스","",
        "rear: 마지막으로 저장된 좌표의 인덱스","",
    ]),

    ("좌표 선형 큐 C 코드", [
        "struct point_queue {",
        "        struct point *points;",
        "        int front;",
        "        int rear;",
        "};",
    ]),

    ("POINT_QUEUE_CREATE(capacity) C 코드", [
        "struct point_queue *point_queue_create(int cap)",
        "{",
        "        struct point_queue *q = ",
        "                malloc(sizeof(struct point_queue));",
        "        if (!q) return NULL;","",
        "        q->points = malloc(sizeof(struct point) * (cap + 1));",
    ]),

    ("POINT_QUEUE_CREATE(capacity) C 코드", [
        "        if (!q->points) {",
        "                free(q);",
        "                return NULL;",
        "        }","",
        "        q->front = 0;",
        "        q->rear = 0;",
        "        return q;",
        "}",
    ]),
    ("POINT_QUEUE_FREE(q) C 코드", [
        "void point_queue_free(struct point_queue *q)",
        "{",
        "        if (!q) return;",
        "        free(q->points);",
        "        free(q);",
        "}",
    ]),
    ("POINT_QUEUE_ENQUEUE(q, point) C 코드", [
        "void point_queue_enqueue(struct point_queue *q, ",
        "                                                          struct point point)",
        "{",
        "        q->points[++q->rear] = point;",
        "}",
    ]),
    ("POINT_QUEUE_DEQUEUE(q) C 코드", [
        "struct point ",
        "point_queue_dequeue(struct point_queue *q)",
        "{",
        "        return q->points[q->front++];",
        "}",
    ]),    

    ("NEXT_POINT(point, direction) 개요", [
        "point의 좌표를 direction 방향으로 이동한 좌표 반환","",
        "distance는 point의 distance에서 1 증가","",
        "direction: 1(상), 2(하), 3(좌), 4(우)",
    ]),

    ("NEXT_POINT(point, direction) C 코드", [
        "struct point next_point(struct point point, int direction)",
        "{",
        "        struct point next = point;",
        "        next.prev_size++;",
        "        if (direction == 1) {			// 상",
        "                next.row--;",
        "        } else if (direction == 2) {		// 하",
        "                next.row++;",
    ]),

    ("NEXT_POINT(point, direction) C 코드", [
        "        } else if (direction == 3) {		// 좌",
        "                next.col--;",
        "        } else if (direction == 4) {		// 우",
        "                next.col++;",
        "        }",
        "        return next;",
        "}",
    ]),

    ("PUSH_NBRS(maze, len, to_visit, point)", [
        "point_index에 point의 row * len + point의 col 저장",
        "반복: point의 상하좌우 좌표 next에 대하여",
        "               만약: STATUS(maze, next) == 0 이면",
        "                           MARK(maze, next, point_index)",
        "                           ENQUEUE(to_visit, next)",
    ]),

    ("PUSH_NBRS(maze, len, to_visit, point)", [
        "void push_nbrs (",
        "       int **maze, ",
        "       int len, ",
        "       struct point_queue *to_visit, ",
        "       struct point point",
        ") {",
        "        int point_index = point.row * len + point.col;", "",
        "        for (int direction = 1; direction <= 4; direction++) {",
    ]),

    ("PUSH_NBRS(maze, len, to_visit, point)", [
        "                struct point next = next_point(point, direction);",
        "                if (maze_status(maze, next) == 0) {",
        "                        maze_mark(maze, next, point_index);",
        "                        point_queue_enqueue(to_visit, next);",
        "                }",
        "        }",
        "}",
    ]),

    ("FIND(maze, len, start, end) 개요 1/2", [
        "to_visit에 CREATE(len * len) 저장     // 갈 곳들",
        "start의 prev_size에 0 저장",
        "MARK(maze, start, -1)",
        "ENQUEUE(to_visit, start)",
        "반복: to_visit이 비어있지 않음",
    ]),

    ("FIND(maze, start, end, len) C 코드 1/2", [
        "int find ("
        "       int **maze, "
        "       int len, "
        "       struct point start, "
        "       struct point end",
        ") {",
        "        struct point_queue *to_visit ",
        "                = point_queue_create(len * len);",
        "        if (!to_visit) return -2;      // 에러코드",
    ]),

    ("FIND(maze, start, end, len) C 코드 1/2", [
        "        start.prev_size = 0;",
        "        maze_mark(maze, start, -1);",
        "        point_queue_enqueue(to_visit, start);","",
        "        while (to_visit->front <= to_visit->rear) {",
    ]),

    ("FIND(maze, len, start, end) 개요 2/2", [
        "               point에 DEQUEUE(to_visit) 저장",
        "               만약: point와 end의 좌표가 동일",
        "                           반환: point.prev_size", 
        "               PUSH_NBRS(maze, len, to_visit, point)","",
        "반환: -1 (못찾음)",
    ]),

    ("FIND(maze, start, end, len) C 코드 2/2", [
        "                struct point point = point_queue_dequeue(to_visit);",
        "                if (point.row == end.row && point.col == end.col) {",
        "                        point_queue_free(to_visit);",
        "                        return point.prev_size; // 도착점까지 거리",
        "                }",
        "                push_nbrs(maze, len, to_visit, point);",    
        "        }",""
        "        point_queue_free(to_visit);",
        "        return -1;",
        "}",
    ]),        

    ("DECODE_PATH(maze, len, end, path_len)", [
        "path에 새 좌표 배열 저장 (크기: path_len + 1)",
        "point에 end 저장",
        "반복: i를 path_len부터 0까지",
        "            path[i]에 point를 저장한다",
        "            prev_data에 STATUS(maze, point)를 저장한다",
        "            point의 row에 prev_data / len를 저장한다",
        r"            point의 col에 prev_data % len를 저장한다",
        "반환: path",
    ]),

    ("DECODE_PATH(maze, len, end, path_len)", [
        "struct point *decode_path("
        "       int **maze, "
        "       int len, "
        "       struct point end, "
        "       int path_len",
        ") {",
        "        struct point *path ",
        "               = malloc(sizeof(struct point) * (path_len + 1));",
        "        if (!path) return NULL;",
    ]),

    ("DECODE_PATH(maze, len, end, path_len)", [
        "        struct point point = end;",
        "        for (int i = path_len; i >= 0; i--) {",
        "                path[i] = point;",
        "                int prev_data = maze_status(maze, point);",
        "                if (prev_data < 0) break;",
        "                point.row = prev_data / len;",
        "                point.col = prev_data % len;",
        "        }",
        "        return path;",
        "}",
    ]),

    ("MAZE_TEST() 개요", [
        "maze 이차원 배열에 미로 정보를 저장한다",
        "m_row_ptrs에 maze의 행 배열들의 배열을 저장한다",
        "출발점, 도착점을 찾고, find()를 실행한다",
        "find() 실행 결과에 따라 에러 처리를 한다",
        "decode_path()를 실행한다",
        "결과를 화면에 출력한다"
    ]),

    ("MAZE_TEST() C 코드", [
        "void maze_test() {",
        "        int maze[7][7] = {",
        "                {1, 1, 1, 1, 1, 1, 1},",
        "                {1, 0, 0, 0, 0, 0, 1},",
        "                {1, 0, 0, 1, 0, 1, 1},",
        "                {1, 0, 1, 0, 0, 0, 1},",
        "                {1, 0, 1, 0, 1, 0, 1},",
        "                {1, 0, 0, 0, 1, 0, 1},",
        "                {1, 1, 1, 1, 1, 1, 1},",
        "        };",
    ]),

    ("MAZE_TEST() C 코드", [
        "        int *m_row_ptrs[7];",
        "        for(int i = 0; i < 7; i++)",
        "                m_row_ptrs[i] = maze[i];","",
        "        struct point start = {1, 1, 0};",
        "        struct point end   = {5, 5, 0};","",
        "        int path_len = find(m_row_ptrs, 7, start, end);",
    ]),

    ("MAZE_TEST() C 코드", [
        "        if (path_len < 0) {",
        '               if (path_len == -1) printf("no path found\\n");',
        '               if (path_len == -2) printf("malloc failed\\n");',
        "               return;",
        "        }","",
        "        struct point *path ",
        "                = decode_path(m_row_ptrs, 7, end, path_len);",
    ]),

    ("MAZE_TEST() C 코드", [
        "        for (int i = 1; i <= path_len; i++)",
        '                printf("(%d,%d), ", path[i].row, path[i].col);',
        '        printf("\\n");',"",
        "       //리소스 정리",
        "       free(path);",
        "}",
    ]),

    ("요약", [
        "단방향 링크드 리스트로 구현한 큐에서 데이터 추가/삭제","",
        "배열로 구현한 큐에서 데이터 추가/삭제","",
        "미로 찾기 예제","",
        "다음주: 계층이라는 차원이 있는 자료구조"
    ]),

    ("C 복습 - 구조체_대입_연산.c", [
        "#include <stdio.h>","",
        "struct Point {",
        "        int x;",
        "        int y;",
        "};","",
        "void struct_test()\n{",
        "        struct Point p1;",
        "        p1.x = 1;",
        "        p1.y = 2;",
        "        struct Point p2 = p1;",
        '        printf("(%d, %d)\\n", p2.x, p2.y);',
        "}","",
        "struct_test()를 실행했을 때 화면에 출력되는 문자열은?"
    ]),

    ("C 복습 - 배열.c", [
        "#include <stdio.h>","",
        "void array_test()\n{",
        "        int arr[3];",
        "        arr[0] = 1;",
        "        arr[1] = 2;",
        "        arr[2] = 3;",
        "        for (int i = 0; i < 3; i++)",
        '                printf("%d ", arr[i]);',"",
        '        printf("\\n");',
        "}","",
        "array_test()를 실행했을 때 화면에 출력되는 문자열은?"
    ]),
    ("C 복습 - 배열_포인터.c", [
        "#include <stdio.h>","",
        "void array_ptr_test()\n{",
        "        int arr[3];",
        "        int *arr_ptr = arr;",
        "        arr_ptr[0] = 1;",
        "        arr_ptr[1] = 2;",
        "        arr_ptr[2] = 3;",
        "        for (int i = 0; i < 3; i++)",
        '                printf("%d ", arr_ptr[i]);',"",
        '        printf("\\n");',
        "}","",
        "array_ptr_test()를 실행했을 때 화면에 출력되는 문자열은?"
    ]),
    ("C 복습 - 배열_포인터_파라미터.c", [
        "#include <stdio.h>","",
        "void array_set_elements(int *arr_ptr)\n{",
        "        arr_ptr[0] = 1;",
        "        arr_ptr[1] = 2;",
        "        arr_ptr[2] = 3;",
        "}","",
        "void array_fun_test()\n{",
        "        int arr[3];",
        "        array_set_elements(arr);",
        "        for (int i = 0; i < 3; i++)",
        '                printf("%d ", arr[i]);',"",
        '        printf("\\n");',
        "}","",
        "array_fun_test()를 실행했을 때 화면에 출력되는 문자열은?"
    ]),

    ("C 코딩 스타일 - if-문의 중괄호", [
        "if와 else에서 모두 문장 하나만 실행할 때는 중괄호를 쓰지 않는다.","",
        'if (condition)\n        do_this();\nelse\n        do_that();'
    ]),
    ("C 코딩 스타일 - if-문의 중괄호 (나쁜 예)", [
        "if의 문장을 숨길 것이 아닌 한 한줄에 쓰지 않는다.","",
        'if (condition) do_this;\ndo_something_everytime;'
    ]),
    ("C 코딩 스타일 - if-문의 중괄호 (나쁜 예)", [
        "괄호를 쓰지 않기 위해서 쉼표를 쓰지 않는다.","",
        'if (condition)\n        do_this(), do_that();'
    ]),
    ("C 코딩 스타일 - if-문의 중괄호", [
        "if나 else 중 하나에만 문장이 두개 이상 있어도 모두 중괄호를 쓴다.","",
        'if (condition) {\n        do_this();\n        do_that();\n} else {\n        otherwise();\n}'
    ]),
    ("C 코딩 스타일 - 반복문의 중괄호", [
        "반복문은 단순 문장 한개만 있지 않은 한 항상 중괄호를 쓴다.","",
        'while (condition) {\n        if (test)\n                do_something();\n}'
    ]),
    ("C 코딩 스타일 - 띄어쓰기 - 키워드", [
        "다음 키워드 다음에는 띄어쓰기를 한칸 한다.","",
        "if, switch, case, for, do, while","",
        "sizeof 다음에는 띄어쓰기를 하지 않는다.","",
        "좋은 예: s = sizeof(struct file);","",
        "괄호 안쪽 주변에는 띄어쓰기를 하지 않는다.","",
        "나쁜 예: s = sizeof( struct file );","",
    ]),
    ("C 코딩 스타일 - 띄어쓰기 - 포인터", [
        "포인터 변수를 선언하거나, 포인터를 반환하는 함수를 선언할 때,","",
        " *은 변수 이름이나 함수 이름에 붙인다.","",
        " *을 타입 옆에 붙이지 않는다.","",
        "char *linux_banner;\nunsigned long long memparse(char *ptr, char **retptr);\nchar *match_strdup(substring_t *s);","",
    ]),
    ("C 코딩 스타일 - 띄어쓰기 - 이항연산자", [
        "다음 이항 연산자 / 삼항 연산자 주변에는 띄어쓰기를 한칸 한다.","",
        "= + -  * / % < > <= >= == != & | ? :","",
        "다음 구조체 멤버 연산자 주변에는 띄어쓰기를 하지 않는다.","",
        ". ->","",
        "area = rect.x  *rect.y;","",
    ]),
    ("C 코딩 스타일 - 띄어쓰기 - 단항연산자", [
        "다음 단항 연산자 뒤에는 띄어쓰기를 하지 않는다.","",
        "&  *+ - ~ ! sizeof","",
        "다음 postfix 증감 단항 연산자 주변에는 띄어쓰기를 하지 않는다.","",
        "++ --","",
        "다음 prefix 증감 단항 연산자 주변에는 띄어쓰기를 하지 않는다.","",
        "++ --","",
        "i++;","",
    ]),
    ("C 코딩 스타일 - 이름", [
        "C 언어는 단순하고 미니멀한 것을 추구한다.","",
        "어떤 프로그래밍 언어에서는 ThisVariableIsATemporaryCounter로 쓰지만","",
        "C 언어에서는 쓰기 훨씬 쉽고, 이해하기는 많이 어렵지 않은 tmp로 쓴다.","",
        "한편, 전역 범위에 있는 이름은 자세하게 만들어야 한다.","",
    ]),
    ("C 코딩 스타일 - 전역 변수 이름, 함수 이름", [
        "전역 변수는 다른 방법이 전혀 없어서 어쩔 수 없이 꼭 필요할 때만 쓴다.","",
        "전역 변수와 함수의 이름은 단어들을 _로 구분해서 자세하게 만든다.","",
        "좋은 예: count_active_users()","",
        "나쁜 예: cntusr()","",
    ]),
    ("C 코딩 스타일 - 지역 변수 이름", [
        "지역 변수 이름은 짧고, 정보를 담고 있어야 한다.","",
        "좋은 예: 루프 변수 i","",
        "나쁜 예: 루프 변수 loop_counter (오해의 소지가 없을 때)","",
        "좋은 예: 임시 변수 tmp","",
    ]),
    ("C 코딩 스타일 - 함수", [
        "함수는 짧고 한가지 일만 해야 한다.","",
        "함수에서 지역 변수를 10개 이하로 써야 한다.","",
        "함수 정의들 간에는 한줄을 띄워서 쓴다.","",
        "함수 선언에서 타입과 이름을 모두 쓴다.","",
    ]),
    ("C 코딩 스타일 - 함수 관련 goto", [
        "goto는 많은 프로그래머들이 더이상 쓰지 않는다.","",
        "여러 군데에서 함수가 종료되는데 공통된 정리 코드가 있을 때 유용하다.","",
        "int fun(int a)\n{\n        int result = 0;\n        char *buffer;\n\n        buffer = malloc(SIZE);\n        if (!buffer)\n                return -ENOMEM;\n\n        if (condition1) {\n                while (loop1) {\n                        ...\n                }\n                result = 1;\n                goto out_free_buffer;\n        }\n        ...\nout_free_buffer:\n        free(buffer);\n        return result;\n}",
    ]),
    ("C 코딩 스타일 - 함수 관련 goto 에러", [
        "one err bugs로 알려진 흔한 버그","",
        "err:\n        free(foo->bar);\n        free(foo);\n        return ret;",
        "foo가 NULL인 상태에서 err로 이동했을 때 버그 발생","",
    ]),
    ("C 코딩 스타일 - 함수 관련 goto 에러 해결책", [
        "err_free_bar:와 err_free_foo:로 라벨을 분리한다.","",
        "err_free_bar:\n        free(foo->bar);\nerr_free_foo:\n        free(foo);\n        return ret;",
        "err_free_bar로는 foo가 NULL이 아닐 때만 이동","",
    ]),
    ("C 코딩 스타일 - 주석", [
        "주석은 함수 앞에 달아서 무엇을 하고, 왜 하는지를 설명한다.","",
        "아주 복잡한 함수의 경우 본문 안에도 구획을 나눠서 주석을 달 수 있다.","",
        "/*\n  *This is the preferred style for multi-line\n  *comments in the Linux kernel source code.\n */","",
        "데이터에 대해 짧은 주석을 다는 것도 중요하다.",""
    ]),
    ("C 복습 - 이름의 범위", [
        "int x = 0;",
        "void function(void)",
        "{",
        "        int a = x;          /* 전역 x(=0) */",
        "        int x = 1;          /* 함수 지역 x */",
        "        for (int x = 2; x < 10; x++) {",
        "                int b = x;  /* for-블록의 x */",
        "        }",
        "        a = x;              /* 여기서의 x는 함수 지역 x(=1) */",
        "}","",
        "함수가 끝나기 전 a에 저장된 값은?",""
    ]),
]
# 슬라이드 생성
for title, bullets in slides:
    add_slide(title, bullets)

# 저장
output_path = "ppt/자료구조_4주차_8회차_PPT.pptx"
prs.save(output_path)
output_path
