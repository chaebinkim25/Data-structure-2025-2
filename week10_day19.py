from pptx import __version__ as pptx_version
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

print(pptx_version)

# 새 프레젠테이션 생성
prs = Presentation("template.pptx")
prs.slides.add_slide(prs.slide_masters[0].slide_layouts[0])  # 제목 슬라이드 추가
# 제목 슬라이드 설정
title_slide = prs.slides[0]
title_slide.shapes.title.text = "자료구조 (Data Structure\n10주차: 그래프"


layout = prs.slide_masters[0].slide_layouts[1]  # 제목과 내용 레이아웃 선택

# 슬라이드 추가 함수
def add_slide(title, bullets):
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    content = slide.placeholders[1]
    content.text = bullets[0] if bullets else ""
    for bullet in bullets[1:]:
        p = content.text_frame.add_paragraph()
        p.text = bullet

slides = [
    ("일차원 자료구조 정리", [
        "리스트: 데이터를 추가하면 순서가 유지되는 구조","",
        "스택: 가장 최신 데이터만 꺼낼 수 있는 구조","",
        "큐: 가장 오래된 데이터만 꺼낼 수 있는 구조","",
    ]),
    ("트리 정리", [
        "트리: 루트로부터 부모 자식 관계로 연결된 구조","",
        "이진 트리: 모든 노드의 자식 노드 숫자가 두개 이하인 트리","",
        "최소힙: 모든 노드의 데이터가 자식 노드 데이터들보다 작은 이진 트리","",
        "힙 정렬: 최소힙으로 정렬되지 않은 구간의 최소값을 반복적으로 꺼내서 정렬하는 방법","",
    ]),
    ("그래프 정리", [
        "그래프: 노드들이 연결선으로 연결된 구조","",
        "깊이 우선 탐색 (DFS): 이웃마다 연결된 모든 노드를 방문해서 트리로 만드는 방법","",
        "너비 우선 탐색 (BFS): 연결선 개수가 짧은 이웃들부터 차례로 트리로 만드는 방법","",
        "최단 경로 (DIJKSTRA): 가까운 이웃부터 차례로 트리로 만드는 방법",
    ]),
    ("이번 시간 목차", [
        "그래프를 최단 길이의 트리로 만드는 방법 (MST)","",
        "최단 길이 트리 알고리즘 예제","",
    ]),
    ("최단 길이 트리", [
        "마을에 전기를 공급하는 비용을 최적화하려면?","",
    ]),
    ("트리보다 더 일반적인 자료 구조", [
        "점들은 선으로 연결. 다대다 관계","",
        "루트가 없기 때문에 ","",
    ]),
    ("그래프란?", [
        "그래프 정의: 각 노드로부터 다른 노드들이 연결된 구조","",
    ]),
    ("그래프", [
        "그래프는 노드와 연결선을 모은 것이다","",
        "\t - vertexes: 전체 노드 리스트 (타입: 그래프 노드)","",
    ]),
    ("그래프 노드 타입", [
        "데이터와 이웃 노드 정보의 묶음","",
        "\t - adj: 이웃 노드 리스트","",
        "\t - value: 저장된 데이터","",
    ]),
    ("그래프 예시", [
        "다음과 같은 그래프를 만들어보자","",
        "노드 A, B, C, D, E를 만든다",
        "LIST_ADD(A->adj, B)",
        "LIST_ADD(A->adj, C)",
        "LIST_ADD(B->adj, C)",
        "LIST_ADD(B->adj, A)",
        "LIST_ADD(C->adj, A)",
        "LIST_ADD(C->adj, B)",
        "LIST_ADD(C->adj, D)",
        "LIST_ADD(D->adj, C)",
        "LIST_ADD(D->adj, E)",
        "LIST_ADD(E->adj, D)",
    ]),

    ("그래프 노드 방문", [
        "목적: 특정 노드와 연결된 모든 노드에 대해 반복하기","",
        "깊이 우선: 한 이웃의 이웃들을 방문하고 다음 이웃 진행",
        "\t 데이터 처리 타이밍:",
        "\t - Preorder: 방문전 처리할 작업",
        "\t - Postorder: 방문 후 처리할 작업",
        "너비 우선: 모든 이웃의 방문이 끝나고 이웃의 이웃 방문",
    ]),

    ("깊이 우선 방문 구현", [
        "그래프 노드 타입에 추가할 정보:","",
        "\t - status : 방문 상태 (방문 전, 중, 후)"
    ]),

    ("깊이 우선 방문 구현 - 사전 작업", [
        "함수 DFS_PREP(그래프 G):",
        "반복: G->vertexes의 모든 노드들 u에 대해,",
        "\tu->status = 방문전",
    ]),

    ("깊이 우선 방문 구현 - 방문 함수", [
        "함수 DFS_VISIT(시작 노드 u):",
        "u->status = 방문중",
        "pre_task(u)",
        "반복: u->adj의 모든 노드들 v에 대해,",
        "\t만약: v->status가 방문전이면,",
        "\t\tDFS_VISIT(v)",
        "u->status = 방문후",
        "post_task(u)"
    ]),

    ("깊이 우선 방문 구현 예제 - 연결된 노드 그룹핑", [
        "그래프 노드 타입에 group 항목 추가 (초기값 0)","",
        "전역 변수: n (초기값 0)","",
        "pre_task(u): u->group = n","",
        "post_task(u): 없음",
    ]),

    ("깊이 우선 방문 구현 예제 - 연결된 노드 그룹핑", [
        "함수 GROUPING_CONNECTED(그래프 G):","",
        "DFS_PREP(G)",
        "n = 0",
        "반복: G의 모든 노드 u에 대해,",
        "\t만약: u->status가 방문전이면,",
        "\t\tn += 1",
        "\t\tDFS_VISIT(G, u)"
    ]),

    ("깊이 우선 방문시 연결 선의 분류", [
        "트리 선: 이웃이 방문전이어서 DFS_VISIT이 실행되는 선",
        "역방향 선: 이웃이 방문중이어서 건너뛰는 선",
        "무방향 그래프에서는 이웃이 방문 후인 경우가 없다",
    ]),

    ("예제: 단절점 찾기", [
        "단절점: 제거하면 그래프가 둘로 나눠지는 노드",
        "예시) 네트워크에서 모든 트래픽이 거쳐가는 서버"
    ]),

    ("예제: 단절점 찾기", [
        "어떤 노드에서 깊이 우선 탐색을 시작했는데, 트리 선이 두개 이상이면 그 노드는 단절점이다",
    ]),

    ("예제: 단절점 찾기", [
        "깊이 우선 탐색 중에, 어떤 노드의 모든 자손에서 그 노드의 조상으로 가는 역방향 선이 없으면, 그 노드는 단절점이다",
    ]),

    ("예제: 단절점 찾기", [
        "그 외의 노드들은 모두 단절점이 아니다",
    ]),

    ("예제: 단절점 찾기", [
        "그래프 노드 타입에 추가할 항목:",
        "\tis_ap: 단절점인지 여부 (초기값 0)",
        "\tvisited_time: 방문된 순서 (초기값 0)",
        "\tlowest_back: 자손에서 역방향 선으로 연결된 가장 먼 조상",
        "전역 변수:",
        "\ttime: 방문 순서를 매기기 위한 전역 시간 (초기값 0)",
    ]),

    ("예제: 단절점 찾기", [
        "함수 FIND_AP(그래프 G):",
        "time = 0",
        "반복: G의 모든 노드 u에 대해,",
        "\tu->is_ap = 0",
        "\tu->visited_time = 0",
        "\tu->lowest_back = 무한대",
        "반복: G의 모든 노드 u에 대해,",
        "\t만약: u->visited_time이 0이면,",
        "\t\tDFS_VISIT_ROOT(u)",
    ]),

    ("예제: 단절점 찾기", [
        "함수 DFS_VISIT_ROOT(시작 노드 u):",
        "time += 1",
        "u->visited_time = time",
        "u->lowest_back = time",
        "root_children_count = 0",
        "반복: u->adj의 모든 노드들 v에 대해,",
        "\t만약: v->visited_time이 0이면,",
        "\t\troot_children_count += 1",
        "\t\tDFS_VISIT_NONROOT(v, u)",
        "만약: root_children_count > 1이면,",
        "\tu->is_ap = 1",
    ]),

    ("예제: 단절점 찾기", [
        "함수 DFS_VISIT_NONROOT(방문할 노드 u, 부모 노드 p):",
        "time += 1",
        "u->visited_time = time",
        "u->lowest_back = time",

        "반복: u->adj의 모든 노드들 v에 대해,",
        "\t만약: v == p면,",
        "\t\t다음 반복으로",
        "\t만약: v->visited_time이 0이 아니면,",
        "\t\t만약: v->visited_time < u->lowest_back이면,",
        "\t\t\tu->lowest_back을 v->visited_time로 업데이트",
        "\t만약: v->visited_time이 0이면,",
        "\t\tDFS_VISIT_NONROOT(v, u)",
        "\t\t만약: v->lowest_back < u->lowest_back이면,",
        "\t\t\tu->lowest_back을 v->lowest_back로 업데이트",
        "\t\t만약: v->lowest_back >= u->visited_time이면,",
        "\t\t\tu->is_ap = 1"
    ]),

    ("요약", [
        "그래프: 노드들이 선으로 연결된 구조","",
        "깊이 우선 탐색: 한 이웃의 이웃들을 먼저 탐색하고, 다른 이웃에 대해 반복","",
        "단절점: 제거하면 그래프가 둘로 나눠지는 노드","",
    ]),


]

# 슬라이드 생성
for title, bullets in slides:
    add_slide(title, bullets)


def add_table_slides(title, bullets, table_data):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title

        content = slide.placeholders[1]
        content.text = bullets[0] if bullets else ""
        for bullet in bullets[1:]:
                p = content.text_frame.add_paragraph()
                p.text = bullet

        cols, rows = len(table_data), len(table_data[0]) if table_data else 0
        left = Inches(0.5)
        top = Inches(5)
        width = Inches(9)
        height = Inches(0.8 + 0.3 * rows)
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        for c in range(cols):
                for r in range(rows):
                        table.cell(r, c).text = table_data[c][r]


# 저장
output_path = "ppt/자료구조_10주차_19회차_PPT.pptx"
prs.save(output_path)
output_path
